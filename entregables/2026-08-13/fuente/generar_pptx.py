# -*- coding: utf-8 -*-
"""
Genera la presentacion de avance quincenal (29-jul -> 13-ago-2026) en .pptx.

Replica el diseno del HTML de esta misma carpeta: marca institucional #ffcc03,
tinta calida, geometria recta y el motivo "de -> a" que atraviesa la quincena.

    python generar_pptx.py

Salida: ../pptx/presentacion_avance_quincenal.pptx
"""
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── Paleta (identica al HTML, tema claro) ──────────────────────────────────
INK        = RGBColor(0x17, 0x16, 0x0F)
INK2       = RGBColor(0x4A, 0x47, 0x3B)
INK3       = RGBColor(0x7C, 0x78, 0x69)
SURFACE    = RGBColor(0xFB, 0xFA, 0xF6)
PANEL      = RGBColor(0xF2, 0xEF, 0xE4)
LINE       = RGBColor(0xDC, 0xD8, 0xC7)
LINE_SOFT  = RGBColor(0xEB, 0xE7, 0xD9)
MARCA      = RGBColor(0xFF, 0xCC, 0x03)
MARCA_INK  = RGBColor(0x7A, 0x5C, 0x00)
OK         = RGBColor(0x1F, 0x6B, 0x3A)
OK_BG      = RGBColor(0xE6, 0xF0, 0xE6)
ALERTA     = RGBColor(0x9C, 0x5A, 0x00)
ALERTA_BG  = RGBColor(0xF7, 0xEC, 0xD7)
GRAVE      = RGBColor(0xA8, 0x1F, 0x1A)
GRAVE_BG   = RGBColor(0xF8, 0xE5, 0xE2)
MARK_ANTES = RGBColor(0xC9, 0xC4, 0xB0)

DISPLAY = "Segoe UI"
BODY    = "Segoe UI"
MONO    = "Consolas"

# ─── Geometria (pulgadas) ───────────────────────────────────────────────────
SW, SH   = 13.333, 7.5
ML       = 0.78
W        = SW - ML * 2
FILETE_H = 0.06
EYEBROW_Y = 0.46
TITULO_Y  = 0.82
CUERPO_Y  = 1.82
PIE_Y     = 6.86
PIE_LINEA = 6.80

TOTAL = 15


# ─── Utilidades ─────────────────────────────────────────────────────────────

def spc(run, centesimas):
    """Espaciado entre caracteres (no expuesto por python-pptx)."""
    run.font._rPr.set('spc', str(int(centesimas)))


def rect(slide, l, t, w, h, color, linea=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if linea is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = linea
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def caja(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


TROZO = re.compile(r'(\*\*.+?\*\*|~.+?~)')


def escribir(p, texto, size, color, font=BODY, bold=False, fuerte=None):
    """Escribe texto con **negrita** y ~acento~ como runs."""
    fuerte = fuerte or INK
    for trozo in TROZO.split(texto):
        if not trozo:
            continue
        r = p.add_run()
        if trozo.startswith('**') and trozo.endswith('**'):
            r.text = trozo[2:-2]
            r.font.bold = True
            r.font.color.rgb = fuerte
        elif trozo.startswith('~') and trozo.endswith('~'):
            r.text = trozo[1:-1]
            r.font.bold = True
            r.font.color.rgb = MARCA_INK
        else:
            r.text = trozo
            r.font.bold = bold
            r.font.color.rgb = color
        r.font.size = Pt(size)
        r.font.name = font
    return p


def parrafo(tf, texto, size, color, font=BODY, bold=False, primero=False,
            space_after=6, line_spacing=1.28, align=None, fuerte=None):
    p = tf.paragraphs[0] if primero else tf.add_paragraph()
    escribir(p, texto, size, color, font, bold, fuerte)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    return p


# ─── Componentes ────────────────────────────────────────────────────────────

def nueva(prs, fondo=SURFACE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = fondo
    rect(s, 0, 0, SW, FILETE_H, MARCA)
    return s


def eyebrow(s, etiqueta, resto):
    tf = caja(s, ML, EYEBROW_Y, W, 0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = etiqueta.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.bold = True
    r.font.color.rgb = MARCA_INK; spc(r, 140)
    r = p.add_run(); r.text = "  ·  " + resto.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 140)


def titulo(s, texto, size=30, y=TITULO_Y, h=0.9):
    tf = caja(s, ML, y, W, h)
    parrafo(tf, texto, size, INK, DISPLAY, bold=True, primero=True,
            space_after=0, line_spacing=1.06)


def pie(s, izq, n):
    rect(s, ML, PIE_LINEA, W, 0.008, LINE)
    tf = caja(s, ML, PIE_Y, W * 0.7, 0.25)
    p = parrafo(tf, izq.upper(), 9, INK3, MONO, primero=True, space_after=0, line_spacing=1)
    for r in p.runs:
        spc(r, 120)
    tf = caja(s, ML + W * 0.7, PIE_Y, W * 0.3, 0.25)
    p = parrafo(tf, "%02d / %d" % (n, TOTAL), 9, INK3, MONO, primero=True,
                space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)
    for r in p.runs:
        spc(r, 120)


def tiles(s, items, y, h=1.42):
    """items: [(cifra, etiqueta), ...]"""
    n = len(items)
    gap = 0.22
    w = (W - gap * (n - 1)) / n
    for i, (cifra, etiqueta) in enumerate(items):
        l = ML + i * (w + gap)
        rect(s, l, y, w, h, PANEL)
        rect(s, l, y, w, 0.055, MARCA)
        tf = caja(s, l + 0.26, y + 0.26, w - 0.5, h - 0.4)
        parrafo(tf, cifra, 28, INK, MONO, bold=True, primero=True,
                space_after=5, line_spacing=1)
        parrafo(tf, etiqueta, 10.5, INK2, BODY, space_after=0, line_spacing=1.22)


def deltas(s, items, y, alto=0.62, gap=0.13):
    """items: [(que, antes, ahora), ...]"""
    for i, (que, antes, ahora) in enumerate(items):
        t = y + i * (alto + gap)
        rect(s, ML, t, W, alto, PANEL)
        rect(s, ML, t, 0.07, alto, MARCA)
        tf = caja(s, ML + 0.32, t + 0.15, W * 0.52, alto - 0.2)
        parrafo(tf, que, 13, INK2, BODY, primero=True, space_after=0, line_spacing=1.1)
        tf = caja(s, ML + W * 0.55, t + 0.16, W * 0.2, alto - 0.2)
        parrafo(tf, antes, 13, INK3, MONO, primero=True, space_after=0,
                line_spacing=1.1, align=PP_ALIGN.RIGHT)
        tf = caja(s, ML + W * 0.755, t + 0.15, 0.3, alto - 0.2)
        parrafo(tf, "→", 13, MARCA_INK, BODY, bold=True, primero=True,
                space_after=0, line_spacing=1.1, align=PP_ALIGN.CENTER)
        tf = caja(s, ML + W * 0.79, t + 0.12, W * 0.21 - 0.1, alto - 0.2)
        parrafo(tf, ahora, 16, INK, MONO, bold=True, primero=True, space_after=0,
                line_spacing=1.05, align=PP_ALIGN.RIGHT)


def panel(s, l, t, w, h, rotulo, parrafos, tipo="neutro", chip=None):
    fondo, barra = PANEL, LINE
    if tipo == "problema":
        fondo, barra = GRAVE_BG, GRAVE
    elif tipo == "solucion":
        fondo, barra = PANEL, MARCA
    rect(s, l, t, w, h, fondo)
    rect(s, l, t, 0.07, h, barra)

    tf = caja(s, l + 0.32, t + 0.24, w - 0.6, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = rotulo.upper()
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 130)
    p.space_after = Pt(0); p.line_spacing = 1
    if chip:
        texto, color, bg = chip
        cw = 0.15 + len(texto) * 0.072
        cl = l + 0.32 + 1.45
        rect(s, cl, t + 0.2, cw, 0.26, bg, linea=color)
        ctf = caja(s, cl, t + 0.245, cw, 0.22)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 80)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1

    tf = caja(s, l + 0.32, t + 0.62, w - 0.62, h - 0.85)
    for i, txt in enumerate(parrafos):
        parrafo(tf, txt, 12.5, INK2, BODY, primero=(i == 0), space_after=9, line_spacing=1.26)


def vineta(p):
    """Vineta cuadrada de marca, alineada por PowerPoint a la primera linea."""
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '182880')
    pPr.set('indent', '-182880')
    sucesores = ('a:buSzTx', 'a:buSzPct', 'a:buSzPts', 'a:buFontTx', 'a:buFont',
                 'a:buNone', 'a:buAutoNum', 'a:buChar', 'a:tabLst', 'a:defRPr', 'a:extLst')
    buClr = pPr.makeelement(qn('a:buClr'), {})
    buClr.append(pPr.makeelement(qn('a:srgbClr'), {'val': 'FFCC03'}))
    pPr.insert_element_before(buClr, *sucesores)
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    pPr.insert_element_before(buFont, 'a:buNone', 'a:buAutoNum', 'a:buChar',
                              'a:tabLst', 'a:defRPr', 'a:extLst')
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    pPr.insert_element_before(buChar, 'a:tabLst', 'a:defRPr', 'a:extLst')


def lista(s, l, t, w, items, size=12.5, gap=0.1):
    tf = caja(s, l, t, w, 3.4)
    for i, txt in enumerate(items):
        p = parrafo(tf, txt, size, INK2, BODY, primero=(i == 0),
                    space_after=13, line_spacing=1.26)
        vineta(p)


def nota(s, texto, y, w=None):
    tf = caja(s, ML, y, w or W, 0.6)
    parrafo(tf, texto, 10.5, INK3, BODY, primero=True, space_after=0, line_spacing=1.3)


def remate(s, texto, y, h=0.9):
    rect(s, ML, y, 0.07, h, MARCA)
    tf = caja(s, ML + 0.32, y + 0.06, W - 0.4, h)
    parrafo(tf, texto, 17, INK, DISPLAY, primero=True, space_after=0, line_spacing=1.22)


def fila(s, y, celdas, alto=0.62, cabecera=False, sub=None, chip=None):
    """celdas: [(texto, ancho_relativo, align), ...] sobre el ancho util."""
    x = ML
    for texto, wr, align in celdas:
        cw = W * wr
        tf = caja(s, x, y + (0.02 if cabecera else 0.06), cw - 0.2, alto)
        if cabecera:
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = texto.upper()
            r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 110)
            p.alignment = align; p.space_after = Pt(0); p.line_spacing = 1
        else:
            parrafo(tf, texto, 12.5, INK2, BODY, primero=True, space_after=0,
                    line_spacing=1.2, align=align)
        x += cw
    if sub:
        tf = caja(s, ML, y + 0.32, W * 0.5, 0.4)
        parrafo(tf, sub, 10, INK3, BODY, primero=True, space_after=0, line_spacing=1.2)
    if chip:
        texto, color, bg = chip
        cw = 0.18 + len(texto) * 0.075
        rect(s, ML + W - cw, y + 0.05, cw, 0.28, bg, linea=color)
        ctf = caja(s, ML + W - cw, y + 0.1, cw, 0.24)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 70)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1
    rect(s, ML, y + alto, W, 0.008, INK if cabecera else LINE_SOFT)


# ─── Las 15 diapositivas ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

# 01 · Portada
s = nueva(prs)
tf = caja(s, ML, 0.85, W, 0.3)
p = parrafo(tf, "Unidad para las Víctimas  ·  PRY-0662064".upper(), 10.5, INK3, MONO,
            primero=True, space_after=0, line_spacing=1)
for r in p.runs:
    spc(r, 160)
tf = caja(s, ML, 2.15, W * 0.92, 2.6)
parrafo(tf, "SICAV Móvil pasó de reconocer 5.000 personas a reconocer 12,68 millones.",
        40, INK, DISPLAY, bold=True, primero=True, space_after=0, line_spacing=1.02)
rect(s, ML, 4.72, 4.0, 0.035, MARCA)
tf = caja(s, ML, 4.92, W, 0.35)
parrafo(tf, "Avance quincenal  ·  29 jul → 13 ago 2026", 14, INK, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML, 6.55, W * 0.6, 0.3)
parrafo(tf, "Sistema de Caracterización de Víctimas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML + W * 0.5, 6.55, W * 0.5, 0.3)
parrafo(tf, "114 cambios versionados en 15 días", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)

# 02 · La quincena en cifras
s = nueva(prs)
eyebrow(s, "Resumen", "La quincena en cifras")
titulo(s, "El sistema dejó de correr contra datos de prueba")
tf = caja(s, ML, 1.75, W * 0.82, 0.5)
parrafo(tf, "Hoy opera contra los datos reales de la Unidad, y la aplicación está publicada "
            "y disponible para el personal de campo.", 14, INK2, BODY, primero=True,
        space_after=0, line_spacing=1.25)
tiles(s, [("12,68 M", "personas que la app reconoce sin conexión"),
          ("5,93 M", "víctimas en el padrón operativo"),
          ("1.158", "encuestadoras cargadas y habilitadas"),
          ("977", "pruebas automáticas en verde")], 2.75)
remate(s, "Una víctima que nunca fue entrevistada ya puede ser atendida aunque no haya señal. "
          "Antes, el sistema respondía que no existía.", 5.0, 0.95)
pie(s, "SICAV Móvil · Avance quincenal", 2)

# 03 · Divisor Avances
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.75, W, 1.6)
parrafo(tf, "01", 78, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.35, W, 0.9)
parrafo(tf, "Avances", 40, INK, DISPLAY, bold=True, primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML, 4.35, W * 0.62, 1.1)
parrafo(tf, "Cuatro frentes cerrados en la quincena: los datos reales en producción, la "
            "operación sin conexión, el puente con el sistema heredado y la capacidad de la "
            "infraestructura.", 14, INK2, BODY, primero=True, space_after=0, line_spacing=1.3)
pie(s, "SICAV Móvil · Avance quincenal", 3)

# 04 · Datos reales en producción
s = nueva(prs)
eyebrow(s, "Avance 1", "Datos reales en producción")
titulo(s, "Dos fuentes cargadas, cruzadas por documento")
y = 1.95
fila(s, y, [("Conjunto", 0.30, PP_ALIGN.LEFT),
            ("Registros", 0.18, PP_ALIGN.RIGHT),
            ("Para qué sirve", 0.52, PP_ALIGN.LEFT)], alto=0.3, cabecera=True)
y += 0.42
for conj, num, uso in [
        ("**Padrón operativo**", "5.926.005", "Personas con caracterización previa"),
        ("**Universo del RUV**", "12.009.492", "Existencia e identidad de toda la población víctima"),
        ("Personas únicas combinadas", "12.677.172", "Cobertura total del sistema"),
        ("Solo en el universo", "8.123.873", "Víctimas reconocidas que nunca fueron entrevistadas")]:
    fila(s, y, [(conj, 0.30, PP_ALIGN.LEFT),
                (num, 0.18, PP_ALIGN.RIGHT),
                (uso, 0.52, PP_ALIGN.LEFT)], alto=0.5)
    y += 0.62
nota(s, "El cruce se hace **por número de documento** y nunca por identificadores internos: "
        "se comprobó que los dos sistemas usan numeraciones distintas, y usarlas habría "
        "vinculado registros de personas diferentes.", y + 0.18)
pie(s, "Avances", 4)

# 05 · Operación sin conexión
s = nueva(prs)
eyebrow(s, "Avance 2", "Operación sin conexión")
titulo(s, "El celular ya conoce a toda la población víctima")
deltas(s, [("Personas reconocidas sin conexión", "5.000", "12.677.172"),
           ("Peso del archivo que descarga el celular", "896 MB", "319 MB"),
           ("Registros con ficha completa", "5.000", "5.001.403")], 1.9)
panel(s, ML, 4.32, W * 0.46, 2.2, "Cambio en pantalla",
      ["Cuando la persona no tiene ficha, la app ahora dice **«Está en el RUV, sin "
       "caracterizar»** y habilita su registro. Antes mostraba un error sin ninguna "
       "acción disponible."])
panel(s, ML + W * 0.5, 4.32, W * 0.5, 2.2, "Publicado",
      ["**APK versión 1.1.0** (compilación 54), distribuida por el código QR institucional "
       "—que no cambia entre versiones— e instalable sobre la anterior sin desinstalar."],
      tipo="solucion")
pie(s, "Avances", 5)

# 06 · Puente con el sistema heredado
s = nueva(prs)
eyebrow(s, "Avance 3", "Puente con el sistema heredado")
titulo(s, "El trabajo viejo y el nuevo dejaron de ser dos mundos")
lista(s, ML, 1.95, W * 0.47, [
    "**1.158 encuestadoras** creadas en SICAV con su identidad del sistema anterior, y cada "
    "una ve en la app lo que hizo en la aplicación vieja.",
    "**111 caracterizaciones completas** que ningún reporte estaba contando, recuperadas y visibles.",
    "**131 caracterizaciones sin autor** reasignadas: seis usuarios tenían la «Ñ» rota en el login."])
lista(s, ML + W * 0.53, 1.95, W * 0.47, [
    "**Hechos victimizantes** leídos directamente del RUV, tras confirmar que son ~tres "
    "catálogos distintos~ y que copiar el número entre ellos habría corrompido 509.442 registros.",
    "**Sincronización de novedades cada 15 minutos**, con interruptor propio.",
    "**Escritura hacia el sistema heredado** completa: los diez pasos, incluidos capítulo y "
    "cierre —los que hacen que el hogar exista para los reportes."])
nota(s, "La escritura usa los procedimientos oficiales del sistema anterior, nunca inserciones "
        "directas: es la única vía que respeta sus validaciones.", 5.75)
pie(s, "Avances", 6)

# 07 · Infraestructura y continuidad
s = nueva(prs)
eyebrow(s, "Avance 4", "Infraestructura y continuidad")
titulo(s, "Capacidad y respaldo, que era lo que faltaba")
deltas(s, [("Espacio disponible para la base de datos", "16 GB", "207 GB"),
           ("Copias de respaldo de la base", "ninguna útil", "3 en rotación")], 1.95)
panel(s, ML, 3.68, W * 0.48, 2.3, "Traslado",
      ["La base se movió al disco ampliado con **12 minutos de interrupción**, verificando "
       "archivo por archivo entre origen y destino antes de reanudar. Ningún otro sistema "
       "del servidor se vio afectado."])
panel(s, ML + W * 0.52, 3.68, W * 0.48, 2.3, "Respaldo semanal",
      ["Copia física automática, **15 GB en 24 minutos**, verificada antes de borrar la más "
       "antigua: un respaldo a medias es peor que ninguno, porque parece uno."])
pie(s, "Avances", 7)

# 08 · Divisor Dificultades
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.6, W, 1.6)
parrafo(tf, "02", 78, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.2, W, 0.9)
parrafo(tf, "Dificultades y soluciones", 40, INK, DISPLAY, bold=True, primero=True,
        space_after=0, line_spacing=1)
tf = caja(s, ML, 4.25, W * 0.62, 1.3)
parrafo(tf, "Cuatro obstáculos reales del periodo. Tres son de tamaño —el del celular, el del "
            "servidor y el del tiempo— y el cuarto es un defecto heredado que no producía "
            "ningún error visible.", 14, INK2, BODY, primero=True, space_after=0, line_spacing=1.3)
pie(s, "SICAV Móvil · Avance quincenal", 8)

# 09 · El tamaño en el celular
s = nueva(prs)
eyebrow(s, "Dificultad 1", "El tamaño en el celular")
titulo(s, "12,7 millones de personas no caben en un teléfono")
panel(s, ML, 1.78, W * 0.44, 2.62, "El problema",
      ["Llevar los datos completos de toda la población víctima al dispositivo habría "
       "superado **1 GB** por encuestadora. El archivo anterior ya pesaba **896 MB** y solo "
       "cubría 5.000 personas.",
       "Descargarlo con datos móviles en terreno no era viable."], tipo="problema")
panel(s, ML + W * 0.48, 1.78, W * 0.52, 2.62, "La solución",
      ["Se separó la pregunta en dos. La ficha completa solo se lleva de quien ya tiene "
       "caracterización. Para el resto basta responder ~«¿esta persona está en el RUV?»~, "
       "y eso cabe en una estructura de **21,7 MB para los 12,7 millones**.",
       "El fundamento es operativo: la persona está frente al encuestador, que le pregunta "
       "sus datos. El dispositivo solo necesita confirmar que procede atenderla."],
      tipo="solucion")
# barras comparativas
for i, (rot, ancho, val, color) in enumerate(
        [("Antes", 1.0, "896 MB", MARK_ANTES), ("Ahora", 0.356, "319 MB", MARCA)]):
    yb = 4.62 + i * 0.5
    tf = caja(s, ML, yb + 0.06, 0.9, 0.3)
    parrafo(tf, rot, 11, INK2, MONO, primero=True, space_after=0, line_spacing=1)
    pista_l, pista_w = ML + 1.0, W - 2.3
    rect(s, pista_l, yb, pista_w, 0.34, LINE_SOFT)
    rect(s, pista_l, yb, pista_w * ancho, 0.34, color)
    tf = caja(s, ML + W - 1.2, yb + 0.05, 1.2, 0.3)
    parrafo(tf, val, 12, INK, MONO, bold=True, primero=True, space_after=0,
            line_spacing=1, align=PP_ALIGN.RIGHT)
nota(s, "Cubriendo 2.535 veces más personas  ·  La estructura no omite a nadie que esté en el "
        "RUV; puede dar una coincidencia falsa en 1 de cada 1.000 consultas de personas ajenas, "
        "que se advierte en pantalla.", 5.78)
pie(s, "Dificultades y soluciones", 9)

# 10 · El tamaño en el servidor
s = nueva(prs)
eyebrow(s, "Dificultad 2", "El tamaño en el servidor")
titulo(s, "El disco no daba para cargar el padrón")
panel(s, ML, 1.82, W * 0.44, 3.3, "El problema",
      ["El servidor tiene un disco de **61 GB compartido con otros sistemas** y llegó a picos "
       "del 87 % de uso. Actualizar 5,9 millones de filas no cabía: la operación proyectaba "
       "**8 horas** y dejaba la base inflada.",
       "Además, la memoria compartida del motor estaba en **64 MB**, con lo cual el "
       "mantenimiento reventaba sin explicación."], tipo="problema")
panel(s, ML + W * 0.48, 1.82, W * 0.52, 3.3, "La solución, en tres tiempos",
      ["**1. Ganar espacio en caliente.** Se retiraron 6 índices sin uso: 2,2 GB liberados y "
       "un 35 % más de velocidad de escritura, sin sacar el sistema de servicio.",
       "**2. Partir el trabajo.** La actualización se dividió en lotes retomables, con "
       "limpieza cada tres. Se cortó dos veces sin daño: el diseño quedó probado en producción.",
       "**3. Mudarse.** La base se trasladó al disco ampliado de la Unidad y pasó de 16 GB a "
       "207 GB disponibles."], tipo="solucion")
nota(s, "Quedan por liberar **33 GB adicionales** cuando se retire la copia del estado anterior, "
        "que se conserva a propósito como vía de retorno durante 14 días.", 5.4)
pie(s, "Dificultades y soluciones", 10)

# 11 · El tiempo de las cargas
s = nueva(prs)
eyebrow(s, "Dificultad 3", "El tiempo de las cargas")
titulo(s, "Procesos de días, reducidos a minutos")
deltas(s, [("Carga inicial del padrón", "42 h", "2 h"),
           ("Aplicación de fechas sobre 3,3 M de registros", "25 h", "minutos"),
           ("Enlace del padrón con el universo", "3 recorridos", "1 recorrido")], 1.88, alto=0.56, gap=0.11)
panel(s, ML, 3.95, W * 0.44, 2.35, "La lección que costó 50 minutos",
      ["Un análisis previo prometía que cierta reescritura sería **20 veces más rápida**. Se "
       "cortó la migración a mitad para aplicarla. Medida de punta a punta, fue un 5 % ~más "
       "lenta~."], tipo="problema")
panel(s, ML + W * 0.48, 3.95, W * 0.52, 2.35, "Por qué, y qué queda anotado",
      ["El límite no era la consulta: la tabla tiene **26 índices**, y cada fila actualizada "
       "escribe en los 26. Son 260 MB por minuto de registro de transacciones que ninguna "
       "consulta evita.",
       "Queda escrito en la migración: **sobre esta tabla, el techo lo ponen los índices**."],
      tipo="solucion")
pie(s, "Dificultades y soluciones", 11)

# 12 · El dato que no era de esa persona
s = nueva(prs)
eyebrow(s, "Dificultad 4", "Un defecto heredado, sin síntoma")
titulo(s, "El sexo, la etnia y la discapacidad eran de otra persona")
panel(s, ML, 1.82, W * 0.44, 3.15, "El problema",
      ["La carga inicial usó como identificador una columna del sistema anterior que **no "
       "identifica a nadie**: es un contador de filas. Los atributos demográficos quedaron "
       "asociados a la persona equivocada.",
       "Coincidencia medida: **50 %, igual que el azar**. No producía ningún error visible; se "
       "detectó al verificar 68 cédulas del territorio, al ver nombres femeninos marcados como "
       "masculinos."],
      tipo="problema", chip=("Riesgo alto", GRAVE, GRAVE_BG))
panel(s, ML + W * 0.48, 1.82, W * 0.52, 3.15, "La solución",
      ["**Contención.** Se verificó que el dato nunca salió hacia las bases de la Unidad, se "
       "bloquearon las vías de propagación y los 5.926.005 registros quedaron marcados "
       "explícitamente como **no verificados**.",
       "**Reemplazo.** El dato ahora se cruza contra el universo del RUV **por número de "
       "documento**. Cobertura medida: **86,1 %**. El resto queda sin dato, que también es "
       "información: «no consta» y «ninguna» no son lo mismo.",
       "**Trazabilidad.** Cada registro guarda de qué fuente salió su estado."], tipo="solucion")
nota(s, "La identidad —documento, nombres, apellidos y fecha de nacimiento— **no está afectada**: "
        "viene de otra fuente. Lo comprometido son atributos que no deciden si una persona puede "
        "ser atendida, pero sí invalidan temporalmente los reportes con enfoque diferencial.", 5.25)
pie(s, "Dificultades y soluciones", 12)

# 13 · Verificación
s = nueva(prs)
eyebrow(s, "Control", "Cómo se comprobó")
titulo(s, "Nada de esto se afirma sin medición")
tiles(s, [("977", "pruebas automáticas en verde — 862 del servidor y 115 de la app"),
          ("68 / 68", "cédulas del territorio reconocidas sin conexión, ninguna invisible"),
          ("33", "de esas cédulas no aparecían antes por ningún medio")], 2.15, h=1.6)
remate(s, "El defecto de datos no lo encontró una revisión programada: lo encontró verificar "
          "68 casos concretos que llegaron desde terreno.", 4.42, 0.95)
nota(s, "Se recomienda mantener esa práctica de forma periódica. Los defectos que no generan "
        "ningún error son los que permanecen más tiempo sin detectarse — este llevaba meses.", 5.8)
pie(s, "SICAV Móvil · Avance quincenal", 13)

# 14 · Lo que sigue
s = nueva(prs)
eyebrow(s, "Siguiente", "Lo que sigue y de qué depende")
titulo(s, "Dos definiciones desbloquean el resto")
y = 1.92
fila(s, y, [("Acción", 0.44, PP_ALIGN.LEFT),
            ("Depende de", 0.34, PP_ALIGN.LEFT),
            ("Estado", 0.22, PP_ALIGN.RIGHT)], alto=0.3, cabecera=True)
y += 0.42
filas = [
    ("**Fuente definitiva del estado en el RUV**",
     "Define quiénes componen el padrón de 5,9 M; no es solo una decisión técnica",
     "Definición de la supervisión del proyecto", ("Pendiente", ALERTA, ALERTA_BG), 0.82),
    ("**Inicio de la operación en campo**",
     "Las credenciales se asignan contra la fecha confirmada; conviene probar un grupo reducido primero",
     "Confirmación de fecha", ("Pendiente", ALERTA, ALERTA_BG), 0.82),
    ("**Recarga del padrón con la fuente definitiva**", None,
     "La decisión anterior", ("Listo para ejecutar", OK, OK_BG), 0.5),
    ("**Ficha completa sin conexión para 4,5 M de personas**", None,
     "Equipo técnico", ("En curso", OK, OK_BG), 0.5),
    ("**Liberación de 33 GB adicionales**", None,
     "14 días de operación normal tras el traslado", ("Programado", OK, OK_BG), 0.5),
]
for accion, sub, dep, chip, alto in filas:
    fila(s, y, [(accion, 0.44, PP_ALIGN.LEFT),
                (dep, 0.34, PP_ALIGN.LEFT),
                ("", 0.22, PP_ALIGN.RIGHT)], alto=alto, sub=sub, chip=chip)
    y += alto + 0.12
pie(s, "SICAV Móvil · Avance quincenal", 14)

# 15 · Cierre
s = nueva(prs)
tf = caja(s, ML, 0.85, W, 0.3)
p = parrafo(tf, "Equipo SICAV Móvil".upper(), 10.5, INK3, MONO, primero=True,
            space_after=0, line_spacing=1)
for r in p.runs:
    spc(r, 160)
tf = caja(s, ML, 1.9, W * 0.9, 1.6)
parrafo(tf, "El sistema ya opera contra los datos reales de la Unidad.", 40, INK, DISPLAY,
        bold=True, primero=True, space_after=0, line_spacing=1.02)
rect(s, ML, 3.62, 4.0, 0.035, MARCA)
tf = caja(s, ML, 3.82, W, 0.35)
parrafo(tf, "Corte: 13 de agosto de 2026", 14, INK, MONO, primero=True,
        space_after=0, line_spacing=1)
lista(s, ML, 4.6, W * 0.47, [
    "**Javier Aguilar** — arquitectura, backend, base de datos, aplicación móvil e infraestructura",
    "**Brando** — frontend web de la plataforma"], size=12)
lista(s, ML + W * 0.53, 4.6, W * 0.47, [
    "**Oscar** — supervisión del proyecto por parte de la Unidad",
    "Proyecto **PRY-0662064** · Unidad para la Atención y Reparación Integral a las Víctimas"], size=12)
tf = caja(s, ML, 6.55, W * 0.6, 0.3)
parrafo(tf, "Cifras medidas contra producción, no estimadas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML + W * 0.5, 6.55, W * 0.5, 0.3)
parrafo(tf, "15 / 15", 10.5, INK3, MONO, primero=True, space_after=0,
        line_spacing=1, align=PP_ALIGN.RIGHT)

destino = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "pptx")
os.makedirs(destino, exist_ok=True)
ruta = os.path.normpath(os.path.join(destino, "presentacion_avance_quincenal.pptx"))
prs.save(ruta)
print("OK:", ruta, "| diapositivas:", len(prs.slides.__iter__.__self__._sldIdLst))
