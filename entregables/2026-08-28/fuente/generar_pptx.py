# -*- coding: utf-8 -*-
"""
Genera la presentacion de avance semanal (14-ago -> 21-ago-2026) en .pptx.

Replica el diseno del HTML de esta misma carpeta: marca institucional #ffcc03,
tinta calida, geometria recta. La estructura es avance / problema / solucion,
con cada problema y su solucion en la MISMA diapositiva.

    python generar_pptx.py

Salida: ../pptx/presentacion_avance_14-21-ago.pptx
"""
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── Paleta (identica al HTML, tema claro) ──────────────────────────────────
INK        = RGBColor(0x17, 0x16, 0x0F)
INK2       = RGBColor(0x4A, 0x47, 0x3B)
INK3       = RGBColor(0x7C, 0x78, 0x69)
SURFACE    = RGBColor(0xFB, 0xFA, 0xF6)
PANEL      = RGBColor(0xF2, 0xEF, 0xE4)
LINE       = RGBColor(0xDC, 0xD8, 0xC7)
LINE_SOFT  = RGBColor(0xEB, 0xE7, 0xD9)
MARCA      = RGBColor(0xFF, 0xCC, 0x03)
MARCA_INK  = RGBColor(0x7A, 0x5C, 0x00)
OK         = RGBColor(0x1F, 0x6B, 0x3A)
OK_BG      = RGBColor(0xE6, 0xF0, 0xE6)
ALERTA     = RGBColor(0x9C, 0x5A, 0x00)
ALERTA_BG  = RGBColor(0xF7, 0xEC, 0xD7)
GRAVE      = RGBColor(0xA8, 0x1F, 0x1A)
GRAVE_BG   = RGBColor(0xF8, 0xE5, 0xE2)
MARK_ANTES = RGBColor(0xC9, 0xC4, 0xB0)

DISPLAY = "Segoe UI"
BODY    = "Segoe UI"
MONO    = "Consolas"

# ─── Geometria (pulgadas) ───────────────────────────────────────────────────
SW, SH   = 13.333, 7.5
ML       = 0.78
W        = SW - ML * 2
FILETE_H = 0.06
EYEBROW_Y = 0.46
TITULO_Y  = 0.82
CUERPO_Y  = 1.82
PIE_Y     = 6.86
PIE_LINEA = 6.80

TOTAL = 15


# ─── Utilidades ─────────────────────────────────────────────────────────────

def spc(run, centesimas):
    """Espaciado entre caracteres (no expuesto por python-pptx)."""
    run.font._rPr.set('spc', str(int(centesimas)))


def rect(slide, l, t, w, h, color, linea=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if linea is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = linea
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def caja(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


TROZO = re.compile(r'(\*\*.+?\*\*|~.+?~)')


def escribir(p, texto, size, color, font=BODY, bold=False, fuerte=None):
    """Escribe texto con **negrita** y ~acento~ como runs."""
    fuerte = fuerte or INK
    for trozo in TROZO.split(texto):
        if not trozo:
            continue
        r = p.add_run()
        if trozo.startswith('**') and trozo.endswith('**'):
            r.text = trozo[2:-2]
            r.font.bold = True
            r.font.color.rgb = fuerte
        elif trozo.startswith('~') and trozo.endswith('~'):
            r.text = trozo[1:-1]
            r.font.bold = True
            r.font.color.rgb = MARCA_INK
        else:
            r.text = trozo
            r.font.bold = bold
            r.font.color.rgb = color
        r.font.size = Pt(size)
        r.font.name = font
    return p


def parrafo(tf, texto, size, color, font=BODY, bold=False, primero=False,
            space_after=6, line_spacing=1.28, align=None, fuerte=None):
    p = tf.paragraphs[0] if primero else tf.add_paragraph()
    escribir(p, texto, size, color, font, bold, fuerte)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    return p


# ─── Componentes ────────────────────────────────────────────────────────────

def nueva(prs, fondo=SURFACE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = fondo
    rect(s, 0, 0, SW, FILETE_H, MARCA)
    return s


def eyebrow(s, etiqueta, resto):
    tf = caja(s, ML, EYEBROW_Y, W, 0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = etiqueta.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.bold = True
    r.font.color.rgb = MARCA_INK; spc(r, 140)
    r = p.add_run(); r.text = "  ·  " + resto.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 140)


def titulo(s, texto, size=30, y=TITULO_Y, h=0.9):
    tf = caja(s, ML, y, W, h)
    parrafo(tf, texto, size, INK, DISPLAY, bold=True, primero=True,
            space_after=0, line_spacing=1.06)


def pie(s, izq, n):
    rect(s, ML, PIE_LINEA, W, 0.008, LINE)
    tf = caja(s, ML, PIE_Y, W * 0.7, 0.25)
    p = parrafo(tf, izq.upper(), 9, INK3, MONO, primero=True, space_after=0, line_spacing=1)
    for r in p.runs:
        spc(r, 120)
    tf = caja(s, ML + W * 0.7, PIE_Y, W * 0.3, 0.25)
    p = parrafo(tf, "%02d / %d" % (n, TOTAL), 9, INK3, MONO, primero=True,
                space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)
    for r in p.runs:
        spc(r, 120)


def tiles(s, items, y, h=1.42):
    """items: [(cifra, etiqueta), ...]"""
    n = len(items)
    gap = 0.22
    w = (W - gap * (n - 1)) / n
    for i, (cifra, etiqueta) in enumerate(items):
        l = ML + i * (w + gap)
        rect(s, l, y, w, h, PANEL)
        rect(s, l, y, w, 0.055, MARCA)
        tf = caja(s, l + 0.26, y + 0.26, w - 0.5, h - 0.4)
        parrafo(tf, cifra, 28, INK, MONO, bold=True, primero=True,
                space_after=5, line_spacing=1)
        parrafo(tf, etiqueta, 10.5, INK2, BODY, space_after=0, line_spacing=1.22)


def deltas(s, items, y, alto=0.62, gap=0.13):
    """items: [(que, antes, ahora), ...]"""
    for i, (que, antes, ahora) in enumerate(items):
        t = y + i * (alto + gap)
        rect(s, ML, t, W, alto, PANEL)
        rect(s, ML, t, 0.07, alto, MARCA)
        tf = caja(s, ML + 0.32, t + 0.15, W * 0.52, alto - 0.2)
        parrafo(tf, que, 13, INK2, BODY, primero=True, space_after=0, line_spacing=1.1)
        tf = caja(s, ML + W * 0.55, t + 0.16, W * 0.2, alto - 0.2)
        parrafo(tf, antes, 13, INK3, MONO, primero=True, space_after=0,
                line_spacing=1.1, align=PP_ALIGN.RIGHT)
        tf = caja(s, ML + W * 0.755, t + 0.15, 0.3, alto - 0.2)
        parrafo(tf, "→", 13, MARCA_INK, BODY, bold=True, primero=True,
                space_after=0, line_spacing=1.1, align=PP_ALIGN.CENTER)
        tf = caja(s, ML + W * 0.79, t + 0.12, W * 0.21 - 0.1, alto - 0.2)
        parrafo(tf, ahora, 16, INK, MONO, bold=True, primero=True, space_after=0,
                line_spacing=1.05, align=PP_ALIGN.RIGHT)


def panel(s, l, t, w, h, rotulo, parrafos, tipo="neutro", chip=None):
    fondo, barra = PANEL, LINE
    if tipo == "problema":
        fondo, barra = GRAVE_BG, GRAVE
    elif tipo == "solucion":
        fondo, barra = PANEL, MARCA
    rect(s, l, t, w, h, fondo)
    rect(s, l, t, 0.07, h, barra)

    tf = caja(s, l + 0.32, t + 0.24, w - 0.6, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = rotulo.upper()
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 130)
    p.space_after = Pt(0); p.line_spacing = 1
    if chip:
        texto, color, bg = chip
        cw = 0.15 + len(texto) * 0.072
        cl = l + 0.32 + 1.45
        rect(s, cl, t + 0.2, cw, 0.26, bg, linea=color)
        ctf = caja(s, cl, t + 0.245, cw, 0.22)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 80)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1

    tf = caja(s, l + 0.32, t + 0.62, w - 0.62, h - 0.85)
    for i, txt in enumerate(parrafos):
        parrafo(tf, txt, 12.5, INK2, BODY, primero=(i == 0), space_after=9, line_spacing=1.26)


def vineta(p):
    """Vineta cuadrada de marca, alineada por PowerPoint a la primera linea."""
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '182880')
    pPr.set('indent', '-182880')
    sucesores = ('a:buSzTx', 'a:buSzPct', 'a:buSzPts', 'a:buFontTx', 'a:buFont',
                 'a:buNone', 'a:buAutoNum', 'a:buChar', 'a:tabLst', 'a:defRPr', 'a:extLst')
    buClr = pPr.makeelement(qn('a:buClr'), {})
    buClr.append(pPr.makeelement(qn('a:srgbClr'), {'val': 'FFCC03'}))
    pPr.insert_element_before(buClr, *sucesores)
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    pPr.insert_element_before(buFont, 'a:buNone', 'a:buAutoNum', 'a:buChar',
                              'a:tabLst', 'a:defRPr', 'a:extLst')
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    pPr.insert_element_before(buChar, 'a:tabLst', 'a:defRPr', 'a:extLst')


def lista(s, l, t, w, items, size=12.5, gap=0.1):
    tf = caja(s, l, t, w, 3.4)
    for i, txt in enumerate(items):
        p = parrafo(tf, txt, size, INK2, BODY, primero=(i == 0),
                    space_after=13, line_spacing=1.26)
        vineta(p)


def nota(s, texto, y, w=None):
    tf = caja(s, ML, y, w or W, 0.6)
    parrafo(tf, texto, 10.5, INK3, BODY, primero=True, space_after=0, line_spacing=1.3)


def remate(s, texto, y, h=0.9):
    rect(s, ML, y, 0.07, h, MARCA)
    tf = caja(s, ML + 0.32, y + 0.06, W - 0.4, h)
    parrafo(tf, texto, 17, INK, DISPLAY, primero=True, space_after=0, line_spacing=1.22)


def fila(s, y, celdas, alto=0.62, cabecera=False, sub=None, chip=None):
    """celdas: [(texto, ancho_relativo, align), ...] sobre el ancho util."""
    x = ML
    for texto, wr, align in celdas:
        cw = W * wr
        tf = caja(s, x, y + (0.02 if cabecera else 0.06), cw - 0.2, alto)
        if cabecera:
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = texto.upper()
            r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 110)
            p.alignment = align; p.space_after = Pt(0); p.line_spacing = 1
        else:
            parrafo(tf, texto, 12.5, INK2, BODY, primero=True, space_after=0,
                    line_spacing=1.2, align=align)
        x += cw
    if sub:
        tf = caja(s, ML, y + 0.32, W * 0.5, 0.4)
        parrafo(tf, sub, 10, INK3, BODY, primero=True, space_after=0, line_spacing=1.2)
    if chip:
        texto, color, bg = chip
        cw = 0.18 + len(texto) * 0.075
        rect(s, ML + W - cw, y + 0.05, cw, 0.28, bg, linea=color)
        ctf = caja(s, ML + W - cw, y + 0.1, cw, 0.24)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 70)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1
    rect(s, ML, y + alto, W, 0.008, INK if cabecera else LINE_SOFT)


# ─── Las 15 diapositivas ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

CHIP_OK    = (OK, OK_BG)
CHIP_PEND  = (ALERTA, ALERTA_BG)
CHIP_GRAVE = (GRAVE, GRAVE_BG)

AV = "Avances"
DS = "Dificultades y soluciones"
SEM = "SICAV Móvil · Avance semanal"

# 01 · Portada
s = nueva(prs)
tf = caja(s, ML, 0.85, W, 0.3)
p = parrafo(tf, "Unidad para las Víctimas  ·  PRY-0662064".upper(), 10.5, INK3, MONO,
            primero=True, space_after=0, line_spacing=1)
for r in p.runs:
    spc(r, 160)
tf = caja(s, ML, 1.95, W * 0.94, 2.6)
parrafo(tf, "Llegaron tres reclamaciones distintas. Las tres tenían una causa distinta "
            "de la que se les atribuía.", 38, INK, DISPLAY, bold=True,
        primero=True, space_after=0, line_spacing=1.03)
rect(s, ML, 4.86, 4.6, 0.035, MARCA)
tf = caja(s, ML, 5.06, W, 0.35)
parrafo(tf, "Avance semanal  ·  21 ago → 28 ago 2026", 14, INK, MONO, primero=True,
        space_after=0, line_spacing=1)
tf = caja(s, ML, 6.55, W * 0.6, 0.3)
parrafo(tf, "SICAV Móvil  ·  Sistema de Caracterización de Víctimas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML + W * 0.5, 6.55, W * 0.5, 0.3)
parrafo(tf, "25 cambios versionados  ·  3 causas raíz halladas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)

# 02 · El periodo en cifras
s = nueva(prs)
eyebrow(s, "Resumen", "El periodo en cifras")
titulo(s, "La semana en que se dejó de suponer y se midió")
tf = caja(s, ML, 1.75, W * 0.88, 0.7)
parrafo(tf, "Se cerraron los doce hallazgos de los informes de calidad v2 y, en paralelo, se "
            "auditó la base de producción para responder tres reclamaciones que llegaron "
            "por correo.", 14, INK2, BODY, primero=True, space_after=0, line_spacing=1.32)
tiles(s, [("12 / 12", "hallazgos de calidad cerrados: 7 de la aplicación y 5 del panel"),
          ("3", "causas raíz identificadas en producción, con evidencia reproducible"),
          ("25", "cambios versionados en el periodo"),
          ("30", "enlaces territoriales con plan de capacitación y anexos listos")],
      2.72, h=1.62)
remate(s, "Ninguna de las tres reclamaciones resultó ser lo que decía el reporte. En los tres "
          "casos la causa apareció al mirar el dato, no al discutir el correo.", 4.72, 1.1)
pie(s, SEM, 2)

# 03 · Divisor avances
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.55, W, 1.5)
parrafo(tf, "01", 76, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.15, W, 1.0)
parrafo(tf, "Avances", 44, INK, DISPLAY, bold=True, primero=True, space_after=0,
        line_spacing=1.02)
tf = caja(s, ML, 4.35, W * 0.72, 1.1)
parrafo(tf, "Tres frentes cerrados: los informes de calidad v2, el plan de capacitación "
            "completo con sus ocho anexos, y lo que quedó corriendo en producción.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.35)
pie(s, SEM, 3)

# 04 · Avance 1 · los informes de calidad v2
s = nueva(prs)
eyebrow(s, "Avance 1", "Informes IGED-QA-C002 y C003 · versión 2")
titulo(s, "Los doce hallazgos, cerrados de punta a punta", size=28)
deltas(s, [("**H-024** · La búsqueda recorría 12 millones de registros sin índice",
            "5,8 s", "2 ms"),
           ("**APK-005** · El avance dividía por obligatorias que las reglas ocultan",
            "0 %", "real"),
           ("**APK-002** · No era intermitente ni de red: eran tres rechazos del servidor",
            "genérico", "del servidor")], 2.05, alto=0.70, gap=0.15)
lista(s, ML, 4.62, W * 0.95,
      ["**H-010 / H-011** — «undefined sesión(es)» y «Página 1 de NaN»: corregido por los dos "
       "lados, backend y panel.",
       "**APK-003** — al abrir el modo sin conexión aparecieron dos defectos más que habrían "
       "hecho engañosa una reprueba.",
       "**Sesiones ya guardadas** — no solo se corrigió el cálculo: se recalcularon las que "
       "estaban en cero."], size=12.5, gap=0.12)
pie(s, AV, 4)

# 05 · Avance 2 · capacitación
s = nueva(prs)
eyebrow(s, "Avance 2", "Capacitación · 1, 3 y 8 de septiembre")
titulo(s, "El plan de capacitación quedó completo, con sus ocho anexos", size=27)
panel(s, ML, 1.92, W * 0.48, 2.65, "El plan",
      ["Tres jornadas de cuatro horas, con dos bloques cada una y práctica guiada sobre el "
       "ambiente real. Grupo A: **16 enlaces**. Grupo B: **14**.",
       "Agenda horaria, temario, metodología, requisitos y los dos listados nominales "
       "completos."],
      tipo="solucion")
panel(s, ML + W * 0.52, 1.92, W * 0.48, 2.65, "Los ocho anexos",
      ["Pre-test y post-test con clave  ·  banco de **32 preguntas** por capítulo  ·  tres "
       "casos de estudio  ·  plantilla de documentación.",
       "Revisión del manual  ·  encuesta de calidad  ·  las siete piezas gráficas  ·  "
       "verificación de dispositivos."],
      tipo="solucion")
nota(s, "Todo anclado al instrumento real: 14 capítulos, 363 preguntas y 276 reglas de salto "
        "del perfil territorial. La revisión del manual arrojó **nueve hallazgos**, cuatro de "
        "prioridad alta y previos a la primera sesión.", 4.95, W * 0.94)
pie(s, AV, 5)

# 06 · Avance 3 · producción
s = nueva(prs)
eyebrow(s, "Avance 3", "Producción")
titulo(s, "Lo que quedó corriendo y verificado")
fila(s, 1.85, [("Componente", 0.2, PP_ALIGN.LEFT), ("Qué cambió", 0.58, PP_ALIGN.LEFT),
               ("Estado", 0.22, PP_ALIGN.LEFT)], alto=0.3, cabecera=True)
FILAS = [("**Aplicación móvil**", "Versión **1.2.3**, con la versión visible en la pantalla de ingreso", "Desplegada", CHIP_OK),
         ("**Panel de Control**", "Trabajo del frente web integrado: paginación y unificación visual de Autorizaciones", "Desplegado", CHIP_OK),
         ("**Autorizaciones**", "Se autoriza a quien está en el RUV y no en el padrón: se le crea la ficha", "En línea", CHIP_OK),
         ("**Backend**", "El despliegue pasa la versión de la aplicación al servidor", "Verificado", CHIP_OK),
         ("**Acceso**", "Carga de claves desde archivo, cifradas con Argon2", "Listo", CHIP_OK)]
for i, (a, b, c, ch) in enumerate(FILAS):
    fila(s, 2.35 + i * 0.80, [(a, 0.2, PP_ALIGN.LEFT), (b, 0.58, PP_ALIGN.LEFT),
                              ("", 0.22, PP_ALIGN.LEFT)], alto=0.72,
         chip=(c, ch[0], ch[1]))
pie(s, AV, 6)

# 07 · Divisor dificultades
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.55, W, 1.5)
parrafo(tf, "02", 76, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.15, W, 1.0)
parrafo(tf, "Dificultades y soluciones", 44, INK, DISPLAY, bold=True, primero=True,
        space_after=0, line_spacing=1.02)
tf = caja(s, ML, 4.45, W * 0.72, 1.1)
parrafo(tf, "Tres reclamaciones llegaron por correo esta semana. Las tres se respondieron "
            "midiendo la base de producción, y en las tres la causa era otra.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.35)
pie(s, SEM, 7)

# 08 · Dificultad 1 · el job en el otro servidor
s = nueva(prs)
eyebrow(s, "Dificultad 1", "Caso 14512 · Las encuestas sin conexión que no llegaban")
titulo(s, "El proceso no estaba apagado. Estaba en el otro servidor.", size=28)
panel(s, ML, 1.92, W * 0.455, 3.42, "El problema",
      ["Once encuestas de Venadillo, cerradas en el dispositivo, nunca llegaron. El informe "
       "externo lo atribuyó a un proceso **deshabilitado desde 2024** y a la falta de un paso "
       "manual que nadie tenía asignado.",
       "Con esa lectura, la conclusión era que había que reactivar un proceso viejo y asignarle "
       "un responsable a una tarea manual."],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.42, "La solución",
      ["El proceso existe **dos veces, con el mismo nombre, en dos servidores**. El del servidor "
       "9 es una copia muerta. El real está en el otro, encendido, y corre todos los días a las "
       "**18:20**.",
       "Funcionó hasta el 14 de agosto. Desde el 16 falla todas las noches en **menos de cuatro "
       "segundos**: perdió el acceso a su carpeta de trabajo tras un reinicio de la base.",
       "No hay que reactivar nada ni asignar a nadie: ~hay que restablecer una carpeta.~"],
      tipo="solucion")
nota(s, "Cuando funcionaba, cada corrida tardaba tres minutos. Ahora muere en uno, y un día en "
        "cero: no alcanza a intentar la conexión. La evidencia estaba en el registro de "
        "ejecuciones desde la primera noche.", 5.55, W * 0.94)
pie(s, DS, 8)

# 09 · Dificultad 2 · la pérdida activa
s = nueva(prs)
eyebrow(s, "Dificultad 2", "Caso 14512 · Lo que está pasando ahora mismo")
titulo(s, "Cada día se borra la captura de ese día")
panel(s, ML, 1.92, W * 0.48, 2.75, "El problema",
      ["La operación copia a diario los archivos a la carpeta de carga, y **al día siguiente los "
       "elimina a mano**. Es la práctica establecida.",
       "Combinada con el proceso que falla desde el 16 de agosto, se está borrando trabajo de "
       "campo **que nunca entró a la base**."],
      tipo="problema")
panel(s, ML + W * 0.52, 1.92, W * 0.48, 2.75, "La solución",
      ["**Suspender hoy la eliminación manual.** No depende de ningún cambio técnico: es una "
       "instrucción a la operación y es la medida más urgente del informe.",
       "Lo ya represado ~sí se recupera~: existe respaldo, y las ventanas de procesamiento se "
       "ampliaron en julio de 30 a **120 días**."],
      tipo="solucion")
remate(s, "No es una pérdida histórica que haya que lamentar: es una pérdida en curso que se "
          "puede detener con un correo.", 5.05, 1.0)
pie(s, DS, 9)

# 10 · Dificultad 3 · la causa que no era
s = nueva(prs)
eyebrow(s, "Dificultad 3", "Segundo informe externo · La transferencia entre servidores")
titulo(s, "La medición era correcta. La causa, no.")
panel(s, ML, 1.92, W * 0.455, 3.42, "El problema",
      ["Un segundo informe reportó una **pérdida silenciosa del 24,5 %** entre servidores, "
       "atribuida a **colisiones de identificador** entre dos secuencias independientes.",
       "De confirmarse, implicaba rediseñar la generación de identificadores en toda la cadena."],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.42, "La solución",
      ["La medición se reprodujo casi exacta. Pero el denominador estaba inflado: las 539 filas "
       "eran **88 archivos** repetidos hasta ocho veces. La magnitud real es de **31 archivos**, "
       "no 131.",
       "Y la causa queda refutada por el propio código: la instrucción de inserción **no incluye "
       "el identificador**, y un disparador lo reasigna con la secuencia local.",
       "De los 11 archivos perdidos, los 11 tenían el identificador ocupado — pero ~de los 77 "
       "que sí llegaron, también los 77.~"],
      tipo="solucion")
nota(s, "El mecanismo real es otro: descarte por nombre repetido, más dos filtros que eliminan "
        "registros sin pasar por el manejo de errores. Ninguno de los dos estaba identificado.",
     5.55, W * 0.94)
pie(s, DS, 10)

# 11 · Dificultad 4 · el tablero GAVE
s = nueva(prs)
eyebrow(s, "Dificultad 4", "Tablero del GAVE · Jornada de Panamá")
titulo(s, "La información sí estaba. Faltaba reconstruir el reporte.", size=28)
panel(s, ML, 1.88, W * 0.44, 2.25, "El problema",
      ["Se requirió la actualización del tablero y un informe explicando **por qué no se efectuó "
       "la carga** de la jornada del 14 al 16 de julio."],
      tipo="problema")
panel(s, ML + W * 0.48, 1.88, W * 0.52, 2.25, "La solución",
      ["Los **51 hogares y 96 registros** estaban cargados, y **cerrados el mismo día de su "
       "captura**. Nunca hubo demora ni pérdida.",
       "El tablero lee una tabla de reporte que **no tiene ningún trabajo programado que la "
       "reconstruya**: se ejecuta a mano."],
      tipo="solucion")
tiles(s, [("51", "hogares de la jornada, en la base desde julio"),
          ("0", "trabajos programados que actualicen el reporte"),
          ("8:54", "hora en que se reconstruyó, el 27 de agosto")], 4.45, h=1.5)
pie(s, DS, 11)

# 12 · Dificultad 5 · la pregunta con un día de plazo
s = nueva(prs)
eyebrow(s, "Dificultad 5", "Inclusión de pregunta en la caracterización")
titulo(s, "Un enunciado en un correo no es una especificación", size=28)
panel(s, ML, 1.92, W * 0.44, 3.42, "El problema",
      ["Se solicitó incluir una pregunta nueva **en todas las versiones** del instrumento, con "
       "soporte normativo, plazo de **un día** y sin pasar por la mesa de servicios."],
      tipo="problema")
panel(s, ML + W * 0.48, 1.92, W * 0.52, 3.42, "La solución",
      ["Se respondió con la medición: «todas las versiones» son **nueve parametrizaciones** — "
       "1.959 preguntas y 1.043 reglas de flujo en producción.",
       "La ubicación pedida no existe: la pregunta de referencia ocupa **cinco posiciones "
       "distintas** y dos niveles de captura según el perfil.",
       "Y hay una dependencia dura: cada opción necesita su código de correspondencia, que "
       "~emite la propia dirección solicitante.~"],
      tipo="solucion")
nota(s, "Estimación entregada: **19 días-persona**, desglosada por actividad, contados desde la "
        "entrega de los insumos. Sin diccionario, flujo, manual y codificación, la variable se "
        "captura y no se puede validar, migrar ni reportar.", 5.55, W * 0.94)
pie(s, DS, 12)

# 13 · Verificación
s = nueva(prs)
eyebrow(s, "Verificación", "Cómo se midió todo esto")
titulo(s, "Todo lo anterior es reproducible")
tf = caja(s, ML, 1.78, W * 0.9, 0.7)
parrafo(tf, "Las tres respuestas se construyeron con consultas de **solo lectura** sobre la base "
            "de producción. No se modificó ningún dato, procedimiento ni trabajo programado.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.32)
lista(s, ML, 2.72, W * 0.95,
      ["**Registro de ejecuciones del proceso** — el estado y el error exacto de cada corrida, "
       "noche por noche, del 28 de julio a hoy.",
       "**Código vigente en producción** — el cuerpo de los procedimientos leído directamente de "
       "la base, no de documentación.",
       "**Cruce entre corridas y filas cargadas** — lo que distingue un proceso que terminó bien "
       "de uno que además hizo su trabajo.",
       "**Recuento sobre las nueve parametrizaciones** — capítulos, preguntas, reglas y opciones, "
       "con y sin código de correspondencia."], size=13, gap=0.16)
nota(s, "Cada informe entregado incluye el anexo con las consultas, para que cualquier tercero "
        "llegue al mismo resultado sin depender de nuestra palabra.", 5.75, W * 0.94)
pie(s, SEM, 13)

# 14 · Lo que sigue
s = nueva(prs)
eyebrow(s, "Lo que sigue", "Decisiones y fechas")
titulo(s, "Dos cosas urgentes, y ninguna es nuestra")
fila(s, 1.85, [("Acción", 0.46, PP_ALIGN.LEFT), ("De quién depende", 0.32, PP_ALIGN.LEFT),
               ("Cuándo", 0.22, PP_ALIGN.LEFT)], alto=0.3, cabecera=True)
SIG = [("**Suspender la eliminación diaria** de archivos del FTP", "Operación", "Hoy", CHIP_GRAVE),
       ("**Restablecer la carpeta de trabajo** en el servidor de Modelo", "Administrador de ese servidor", "Inmediato", CHIP_GRAVE),
       ("Reprocesar el respaldo de los archivos represados", "Equipo técnico, tras lo anterior", "Antes de octubre", CHIP_PEND),
       ("Programar la reconstrucción del reporte del GAVE", "Autorización de la Subdirección", "Septiembre", CHIP_PEND),
       ("Entrega de insumos para la pregunta nueva", "Dirección de Registro", "Sin fecha", CHIP_PEND),
       ("**Capacitación** a los 30 enlaces territoriales", "Equipo · confirmado", "1, 3 y 8 sept", CHIP_OK)]
for i, (a, b, c, ch) in enumerate(SIG):
    fila(s, 2.32 + i * 0.62, [(a, 0.46, PP_ALIGN.LEFT), (b, 0.32, PP_ALIGN.LEFT),
                              ("", 0.22, PP_ALIGN.LEFT)], alto=0.56,
         chip=(c, ch[0], ch[1]))
nota(s, "La ventana de recuperación se cierra sola: para capturas de mediados de julio, el plazo "
        "vence alrededor del **12 de octubre**.", 6.12, W * 0.94)
pie(s, SEM, 14)

# 15 · Cierre
s = nueva(prs)
eyebrow(s, "Cierre", "21 ago → 28 ago 2026")
titulo(s, "Los doce hallazgos, cerrados. Y tres causas raíz que nadie había ubicado.", size=27)
tiles(s, [("12 / 12", "hallazgos de calidad cerrados"),
          ("3", "causas raíz halladas, con evidencia reproducible"),
          ("4", "informes técnicos entregados"),
          ("30", "enlaces con capacitación lista")], 2.15, h=1.5)
remate(s, "La diferencia entre un reporte y un diagnóstico es una consulta. Esta semana se "
          "hicieron las consultas.", 4.02, 1.0)
nota(s, "Javier Aguilar · Arquitectura y desarrollo    ·    Brandon · Panel de Control    ·    "
        "Jorge · Calidad    ·    PRY-0662064 · Contrato 2226-2026", 5.55, W * 0.94)
pie(s, SEM, 15)

# ─── Guardar ────────────────────────────────────────────────────────────────
import os
DEST = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'pptx',
                    'presentacion_avance_21-28-ago.pptx')
prs.save(os.path.normpath(DEST))
print("PPTX generado:", os.path.normpath(DEST))
print("diapositivas:", len(prs.slides.__iter__.__self__._sldIdLst))
