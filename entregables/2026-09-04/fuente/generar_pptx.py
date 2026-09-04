# -*- coding: utf-8 -*-
"""
Genera la presentacion de avance semanal (28-ago -> 4-sep-2026) en .pptx.

Replica el diseno del HTML de esta misma carpeta: marca institucional #ffcc03,
tinta calida, geometria recta. La estructura es avance / dificultad / solucion,
con cada dificultad y su solucion en la MISMA diapositiva.

    python generar_pptx.py

Salida: ../pptx/presentacion_avance_28ago-4sep.pptx
"""
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── Paleta (identica al HTML, tema claro) ──────────────────────────────────
INK        = RGBColor(0x17, 0x16, 0x0F)
INK2       = RGBColor(0x4A, 0x47, 0x3B)
INK3       = RGBColor(0x7C, 0x78, 0x69)
SURFACE    = RGBColor(0xFB, 0xFA, 0xF6)
PANEL      = RGBColor(0xF2, 0xEF, 0xE4)
LINE       = RGBColor(0xDC, 0xD8, 0xC7)
LINE_SOFT  = RGBColor(0xEB, 0xE7, 0xD9)
MARCA      = RGBColor(0xFF, 0xCC, 0x03)
MARCA_INK  = RGBColor(0x7A, 0x5C, 0x00)
OK         = RGBColor(0x1F, 0x6B, 0x3A)
OK_BG      = RGBColor(0xE6, 0xF0, 0xE6)
ALERTA     = RGBColor(0x9C, 0x5A, 0x00)
ALERTA_BG  = RGBColor(0xF7, 0xEC, 0xD7)
GRAVE      = RGBColor(0xA8, 0x1F, 0x1A)
GRAVE_BG   = RGBColor(0xF8, 0xE5, 0xE2)
MARK_ANTES = RGBColor(0xC9, 0xC4, 0xB0)

DISPLAY = "Segoe UI"
BODY    = "Segoe UI"
MONO    = "Consolas"

# ─── Geometria (pulgadas) ───────────────────────────────────────────────────
SW, SH   = 13.333, 7.5
ML       = 0.78
W        = SW - ML * 2
FILETE_H = 0.06
EYEBROW_Y = 0.46
TITULO_Y  = 0.82
CUERPO_Y  = 1.82
PIE_Y     = 6.86
PIE_LINEA = 6.80

TOTAL = 15


# ─── Utilidades ─────────────────────────────────────────────────────────────

def spc(run, centesimas):
    """Espaciado entre caracteres (no expuesto por python-pptx)."""
    run.font._rPr.set('spc', str(int(centesimas)))


def rect(slide, l, t, w, h, color, linea=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if linea is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = linea
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def caja(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


TROZO = re.compile(r'(\*\*.+?\*\*|~.+?~)')


def escribir(p, texto, size, color, font=BODY, bold=False, fuerte=None):
    """Escribe texto con **negrita** y ~acento~ como runs."""
    fuerte = fuerte or INK
    for trozo in TROZO.split(texto):
        if not trozo:
            continue
        r = p.add_run()
        if trozo.startswith('**') and trozo.endswith('**'):
            r.text = trozo[2:-2]
            r.font.bold = True
            r.font.color.rgb = fuerte
        elif trozo.startswith('~') and trozo.endswith('~'):
            r.text = trozo[1:-1]
            r.font.bold = True
            r.font.color.rgb = MARCA_INK
        else:
            r.text = trozo
            r.font.bold = bold
            r.font.color.rgb = color
        r.font.size = Pt(size)
        r.font.name = font
    return p


def parrafo(tf, texto, size, color, font=BODY, bold=False, primero=False,
            space_after=6, line_spacing=1.28, align=None, fuerte=None):
    p = tf.paragraphs[0] if primero else tf.add_paragraph()
    escribir(p, texto, size, color, font, bold, fuerte)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    return p


# ─── Componentes ────────────────────────────────────────────────────────────

def nueva(prs, fondo=SURFACE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = fondo
    rect(s, 0, 0, SW, FILETE_H, MARCA)
    return s


def eyebrow(s, etiqueta, resto):
    tf = caja(s, ML, EYEBROW_Y, W, 0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = etiqueta.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.bold = True
    r.font.color.rgb = MARCA_INK; spc(r, 140)
    r = p.add_run(); r.text = "  ·  " + resto.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 140)


def titulo(s, texto, size=30, y=TITULO_Y, h=0.9):
    tf = caja(s, ML, y, W, h)
    parrafo(tf, texto, size, INK, DISPLAY, bold=True, primero=True,
            space_after=0, line_spacing=1.06)


def pie(s, izq, n):
    rect(s, ML, PIE_LINEA, W, 0.008, LINE)
    tf = caja(s, ML, PIE_Y, W * 0.7, 0.25)
    p = parrafo(tf, izq.upper(), 9, INK3, MONO, primero=True, space_after=0, line_spacing=1)
    for r in p.runs:
        spc(r, 120)
    tf = caja(s, ML + W * 0.7, PIE_Y, W * 0.3, 0.25)
    p = parrafo(tf, "%02d / %d" % (n, TOTAL), 9, INK3, MONO, primero=True,
                space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)
    for r in p.runs:
        spc(r, 120)


def tiles(s, items, y, h=1.42):
    """items: [(cifra, etiqueta), ...]"""
    n = len(items)
    gap = 0.22
    w = (W - gap * (n - 1)) / n
    for i, (cifra, etiqueta) in enumerate(items):
        l = ML + i * (w + gap)
        rect(s, l, y, w, h, PANEL)
        rect(s, l, y, w, 0.055, MARCA)
        tf = caja(s, l + 0.26, y + 0.26, w - 0.5, h - 0.4)
        parrafo(tf, cifra, 28, INK, MONO, bold=True, primero=True,
                space_after=5, line_spacing=1)
        parrafo(tf, etiqueta, 10.5, INK2, BODY, space_after=0, line_spacing=1.22)


def deltas(s, items, y, alto=0.62, gap=0.13):
    """items: [(que, antes, ahora), ...]"""
    for i, (que, antes, ahora) in enumerate(items):
        t = y + i * (alto + gap)
        rect(s, ML, t, W, alto, PANEL)
        rect(s, ML, t, 0.07, alto, MARCA)
        tf = caja(s, ML + 0.32, t + 0.15, W * 0.52, alto - 0.2)
        parrafo(tf, que, 13, INK2, BODY, primero=True, space_after=0, line_spacing=1.1)
        tf = caja(s, ML + W * 0.55, t + 0.16, W * 0.2, alto - 0.2)
        parrafo(tf, antes, 13, INK3, MONO, primero=True, space_after=0,
                line_spacing=1.1, align=PP_ALIGN.RIGHT)
        tf = caja(s, ML + W * 0.755, t + 0.15, 0.3, alto - 0.2)
        parrafo(tf, "→", 13, MARCA_INK, BODY, bold=True, primero=True,
                space_after=0, line_spacing=1.1, align=PP_ALIGN.CENTER)
        tf = caja(s, ML + W * 0.79, t + 0.12, W * 0.21 - 0.1, alto - 0.2)
        parrafo(tf, ahora, 16, INK, MONO, bold=True, primero=True, space_after=0,
                line_spacing=1.05, align=PP_ALIGN.RIGHT)


def panel(s, l, t, w, h, rotulo, parrafos, tipo="neutro", chip=None):
    fondo, barra = PANEL, LINE
    if tipo == "problema":
        fondo, barra = GRAVE_BG, GRAVE
    elif tipo == "solucion":
        fondo, barra = PANEL, MARCA
    rect(s, l, t, w, h, fondo)
    rect(s, l, t, 0.07, h, barra)

    tf = caja(s, l + 0.32, t + 0.24, w - 0.6, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = rotulo.upper()
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 130)
    p.space_after = Pt(0); p.line_spacing = 1
    if chip:
        texto, color, bg = chip
        cw = 0.15 + len(texto) * 0.072
        cl = l + 0.32 + 1.45
        rect(s, cl, t + 0.2, cw, 0.26, bg, linea=color)
        ctf = caja(s, cl, t + 0.245, cw, 0.22)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 80)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1

    tf = caja(s, l + 0.32, t + 0.62, w - 0.62, h - 0.85)
    for i, txt in enumerate(parrafos):
        parrafo(tf, txt, 12.5, INK2, BODY, primero=(i == 0), space_after=9, line_spacing=1.26)


def vineta(p):
    """Vineta cuadrada de marca, alineada por PowerPoint a la primera linea."""
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '182880')
    pPr.set('indent', '-182880')
    sucesores = ('a:buSzTx', 'a:buSzPct', 'a:buSzPts', 'a:buFontTx', 'a:buFont',
                 'a:buNone', 'a:buAutoNum', 'a:buChar', 'a:tabLst', 'a:defRPr', 'a:extLst')
    buClr = pPr.makeelement(qn('a:buClr'), {})
    buClr.append(pPr.makeelement(qn('a:srgbClr'), {'val': 'FFCC03'}))
    pPr.insert_element_before(buClr, *sucesores)
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    pPr.insert_element_before(buFont, 'a:buNone', 'a:buAutoNum', 'a:buChar',
                              'a:tabLst', 'a:defRPr', 'a:extLst')
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    pPr.insert_element_before(buChar, 'a:tabLst', 'a:defRPr', 'a:extLst')


def lista(s, l, t, w, items, size=12.5, gap=0.1):
    tf = caja(s, l, t, w, 3.4)
    for i, txt in enumerate(items):
        p = parrafo(tf, txt, size, INK2, BODY, primero=(i == 0),
                    space_after=13, line_spacing=1.26)
        vineta(p)


def nota(s, texto, y, w=None):
    tf = caja(s, ML, y, w or W, 0.6)
    parrafo(tf, texto, 10.5, INK3, BODY, primero=True, space_after=0, line_spacing=1.3)


def remate(s, texto, y, h=0.9):
    rect(s, ML, y, 0.07, h, MARCA)
    tf = caja(s, ML + 0.32, y + 0.06, W - 0.4, h)
    parrafo(tf, texto, 17, INK, DISPLAY, primero=True, space_after=0, line_spacing=1.22)


def fila(s, y, celdas, alto=0.62, cabecera=False, sub=None, chip=None):
    """celdas: [(texto, ancho_relativo, align), ...] sobre el ancho util."""
    x = ML
    for texto, wr, align in celdas:
        cw = W * wr
        tf = caja(s, x, y + (0.02 if cabecera else 0.06), cw - 0.2, alto)
        if cabecera:
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = texto.upper()
            r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 110)
            p.alignment = align; p.space_after = Pt(0); p.line_spacing = 1
        else:
            parrafo(tf, texto, 12.5, INK2, BODY, primero=True, space_after=0,
                    line_spacing=1.2, align=align)
        x += cw
    if sub:
        tf = caja(s, ML, y + 0.32, W * 0.5, 0.4)
        parrafo(tf, sub, 10, INK3, BODY, primero=True, space_after=0, line_spacing=1.2)
    if chip:
        texto, color, bg = chip
        cw = 0.18 + len(texto) * 0.075
        rect(s, ML + W - cw, y + 0.05, cw, 0.28, bg, linea=color)
        ctf = caja(s, ML + W - cw, y + 0.1, cw, 0.24)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 70)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1
    rect(s, ML, y + alto, W, 0.008, INK if cabecera else LINE_SOFT)



# ─── Las 15 diapositivas ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

CHIP_OK    = (OK, OK_BG)
CHIP_PEND  = (ALERTA, ALERTA_BG)
CHIP_GRAVE = (GRAVE, GRAVE_BG)

AV = "Avances"
DS = "Dificultades y soluciones"
SEM = "SICAV Móvil · Avance semanal"

# 01 · Portada
s = nueva(prs)
tf = caja(s, ML, 0.85, W, 0.3)
p = parrafo(tf, "Unidad para las Víctimas  ·  PRY-0662064".upper(), 10.5, INK3, MONO,
            primero=True, space_after=0, line_spacing=1)
for r in p.runs:
    spc(r, 160)
tf = caja(s, ML, 1.95, W * 0.94, 2.6)
parrafo(tf, "La capacitación dejó de ser un acto que se dicta y pasó a ser algo que se mide.",
        38, INK, DISPLAY, bold=True, primero=True, space_after=0, line_spacing=1.03)
rect(s, ML, 4.86, 4.6, 0.035, MARCA)
tf = caja(s, ML, 5.06, W, 0.35)
parrafo(tf, "Avance semanal  ·  28 ago → 4 sep 2026", 14, INK, MONO, primero=True,
        space_after=0, line_spacing=1)
tf = caja(s, ML, 6.55, W * 0.6, 0.3)
parrafo(tf, "SICAV Móvil  ·  Sistema de Caracterización de Víctimas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML + W * 0.5, 6.55, W * 0.5, 0.3)
parrafo(tf, "14 cambios versionados  ·  1.185 pruebas en verde", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)

# 02 · El periodo en cifras
s = nueva(prs)
eyebrow(s, "Resumen", "El periodo en cifras")
titulo(s, "Se cerró el alistamiento de la capacitación y se auditó el proyecto", size=28)
tf = caja(s, ML, 1.75, W * 0.9, 0.7)
parrafo(tf, "La semana tuvo dos mitades. En la primera se midió el estado real de cada frente "
            "—incluidos cinco pendientes que ya no lo eran—. En la segunda se construyó lo que "
            "faltaba para capacitar: un cuestionario que califica solo, el manual publicado y "
            "la aplicación descargable.", 13.5, INK2, BODY, primero=True, space_after=0,
        line_spacing=1.3)
tiles(s, [("1.185", "pruebas automáticas en verde, 0 fallos: 1.037 del backend y 148 de la aplicación"),
          ("10", "preguntas del pre-test: cinco minutos, publicadas y calificando en el servidor"),
          ("61", "comprobaciones de permisos que antes eran una lista de chequeo a mano"),
          ("4", "correos pendientes enviados el 1 de septiembre")],
      2.86, h=1.62)
remate(s, "Cuatro de las cinco dificultades de la semana las encontramos nosotros revisando "
          "nuestro propio trabajo, no un informe externo.", 4.85, 1.0)
pie(s, SEM, 2)

# 03 · Divisor avances
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.55, W, 1.5)
parrafo(tf, "01", 76, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.15, W, 1.0)
parrafo(tf, "Avances", 44, INK, DISPLAY, bold=True, primero=True, space_after=0,
        line_spacing=1.02)
tf = caja(s, ML, 4.35, W * 0.72, 1.1)
parrafo(tf, "Tres frentes: el instrumento para medir la capacitación, la auditoría del estado "
            "real del proyecto, y lo que quedó corriendo en producción.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.35)
pie(s, SEM, 3)

# 04 · Avance 1 · el pre-test y el post-test
s = nueva(prs)
eyebrow(s, "Avance 1", "Capacitación · Medir la jornada, no solo dictarla")
titulo(s, "Un cuestionario que califica solo y dice cuánto se aprendió", size=28)
tf = caja(s, ML, 1.78, W * 0.9, 0.6)
parrafo(tf, "Diez preguntas —cinco minutos—, el **mismo cuestionario** al inicio y al cierre. Lo que importa "
            "no es el puntaje final aislado: es la diferencia entre los dos, persona por persona.",
        13.5, INK2, BODY, primero=True, space_after=0, line_spacing=1.3)
lista(s, ML, 2.62, W * 0.95,
      ["**Se responde en línea**, sin usuario ni contraseña. Cada quien se identifica con su "
       "correo institucional, que es lo que empareja el pre con el post.",
       "**Califica el servidor, no el navegador.** La clave de respuestas nunca viaja al "
       "dispositivo, así que el cuestionario no se puede leer por adelantado.",
       "**Devuelve el resultado al instante**, con la explicación de cada error: el participante "
       "aprende en el momento en que se equivoca.",
       "**Las preguntas salen del sistema real** —la regla de los dos años, quién autoriza la "
       "excepción, los 14 capítulos, el indicador «✓ Al día»—, no de un temario genérico."],
      size=12.5, gap=0.12)
remate(s, "Sin esto, de una capacitación solo queda la lista de asistencia. Con esto queda una "
          "medición de cuánto subió cada persona.", 5.45, 0.95)
pie(s, AV, 4)

# 05 · Avance 2 · el estado global medido
s = nueva(prs)
eyebrow(s, "Avance 2", "Auditoría interna · El estado real de cada frente")
titulo(s, "Se corrió todo y se revisó todo, incluidos los pendientes", size=28)
deltas(s, [("**Matriz de permisos** · era una lista de chequeo a mano, y nunca se diligenció",
            "manual", "61 pruebas"),
           ("**Batería del backend** · antes cubría tres perfiles; ahora los cinco reales",
            "976", "1.037")], 1.95, alto=0.70, gap=0.15)
lista(s, ML, 3.62, W * 0.95,
      ["**Cinco pendientes cerrados por comprobación**, no por olvido: entre ellos el «bug del "
       "403 para el Supervisor», que ya no existía, y un ajuste de datos que resultó no cambiar "
       "nada en producción.",
       "**Un hallazgo colateral:** el perfil COORDINADOR puede caracterizar ~y~ autorizar "
       "excepciones a la vez, lo que contradice el criterio con que se creó ese permiso.",
       "**Medido, no supuesto:** hay **3 cuentas** con permiso de autorizar para "
       "**1.157 encuestadores**."], size=12.5, gap=0.12)
pie(s, AV, 5)

# 06 · Avance 3 · lo que quedó corriendo
s = nueva(prs)
eyebrow(s, "Avance 3", "Producción · Lo que quedó publicado y funcionando")
titulo(s, "Todo lo que los treinta enlaces necesitan, ya está en línea", size=28)
fila(s, 1.85, [("Qué", 0.56, PP_ALIGN.LEFT), ("Para quién", 0.44, PP_ALIGN.LEFT)],
     alto=0.3, cabecera=True)
FILAS = [("**Aplicación SICAV Móvil 1.2.3** — con código QR y pasos de instalación",
          "Los 30 enlaces territoriales", "Publicada", CHIP_OK),
         ("**Manual de Uso v1.2** — para encuestadores",
          "Encuestadores y enlaces", "Publicado", CHIP_OK),
         ("**Pre-test** — 10 preguntas, 5 minutos, califica en el servidor",
          "Antes de cada jornada", "Abierto", CHIP_OK),
         ("**Post-test** — el mismo cuestionario",
          "Al cierre de cada jornada", "Abierto", CHIP_OK)]
for i, (a, b, c, ch) in enumerate(FILAS):
    fila(s, 2.35 + i * 0.72, [(a, 0.56, PP_ALIGN.LEFT), (b, 0.28, PP_ALIGN.LEFT),
                              ("", 0.16, PP_ALIGN.LEFT)], alto=0.64,
         chip=(c, ch[0], ch[1]))
nota(s, "El Manual de Uso pasó a la versión 1.2 con los cuatro hallazgos de prioridad alta de la "
        "revisión: la regla de los dos años y cómo se levanta, cómo confirmar la versión "
        "instalada, cómo reportar un problema sin incluir datos de la persona entrevistada, y "
        "tres casos nuevos en la guía de problemas.", 5.45, W * 0.94)
pie(s, AV, 6)

# 07 · Divisor dificultades
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.55, W, 1.5)
parrafo(tf, "02", 76, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.15, W, 1.0)
parrafo(tf, "Dificultades y soluciones", 44, INK, DISPLAY, bold=True, primero=True,
        space_after=0, line_spacing=1.02)
tf = caja(s, ML, 4.45, W * 0.72, 1.1)
parrafo(tf, "Cinco. Cuatro salieron de revisar nuestro propio trabajo antes de que llegaran a "
            "nadie. La quinta lleva veinte noches y no depende de nosotros.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.35)
pie(s, SEM, 7)

# 08 · Dificultad 1 · el cuestionario no medía nada
s = nueva(prs)
eyebrow(s, "Dificultad 1", "Pre-test · Encontrada antes de aplicarla a nadie")
titulo(s, "El cuestionario no medía nada: bastaba marcar «todo B»", size=28)
panel(s, ML, 1.92, W * 0.48, 3.15, "El problema",
      ["En la primera versión, **once de las quince respuestas correctas eran la B** y ninguna "
       "era la D. Quien marcara «todo B» sin leer sacaba **11 sobre 15**: nivel «suficiente», "
       "habilitado para operar.",
       "Corregida la clave, apareció la misma fuga por otra vía: la correcta era **la opción más "
       "larga en 7 de cada 10**. Marcar «la más larga» seguía aprobando."],
      tipo="problema")
panel(s, ML + W * 0.52, 1.92, W * 0.48, 3.15, "La solución",
      ["Clave repartida entre las cuatro letras y distractores alargados hasta que **la correcta "
       "no es la más larga en ninguna**. Hoy la letra que más se repite da 3 sobre 10: responder "
       "a ciegas no aprueba.",
       "Los rótulos se verificaron **contra el código de la aplicación** —«Registrar y "
       "caracterizar», «✓ Al día»— y no contra los nombres internos con que hablamos entre "
       "nosotros.",
       "Y por pedido de la Subdirección quedó en **10 preguntas, cinco minutos**."],
      tipo="solucion")
remate(s, "Un examen mal armado no es inofensivo: habría certificado como habilitados a personas "
          "que no lo están, y con constancia escrita.", 5.35, 0.95)
pie(s, DS, 8)

# 09 · Dificultad 2 · el manual no estaba publicado
s = nueva(prs)
eyebrow(s, "Dificultad 2", "Manual de Uso · Hallada al revisar el anexo")
titulo(s, "El manual existía, estaba bien escrito y no estaba en ninguna parte", size=27)
panel(s, ML, 1.92, W * 0.455, 3.25, "El problema",
      ["La revisión del manual daba por hecho que estaba publicado. No lo estaba: la página de "
       "descarga solo ofrecía la aplicación.",
       "Lo que sí se servía como «manual» era el **Manual Funcional**, escrito para perfiles de "
       "prueba y soporte. **Los treinta enlaces no tenían de dónde bajar el suyo.**"],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.25, "La solución",
      ["Se publicó el **Manual de Uso v1.2** como página propia, separado y claramente rotulado "
       "frente al Funcional, en la misma pantalla donde se descarga la aplicación.",
       "De paso, la revisión corrigió un hallazgo suyo propio: se había anotado que «el manual no "
       "cubre el alta manual», y sí la cubre —bajo otro nombre—. ~El registro quedó con esa "
       "corrección, no reescrito.~"],
      tipo="solucion")
nota(s, "La revisión se hizo para encontrar defectos en el contenido del manual y terminó "
        "encontrando uno de distribución, que era más grave: un manual impecable que nadie puede "
        "abrir no cumple ninguna función.", 5.42, W * 0.94)
pie(s, DS, 9)

# 10 · Dificultad 3 · la dirección de descarga
s = nueva(prs)
eyebrow(s, "Dificultad 3", "Distribución de la aplicación · Detrás del cortafuegos")
titulo(s, "La dirección para descargar la app dependía de quién preguntara", size=27)
panel(s, ML, 1.92, W * 0.505, 3.15, "El problema",
      ["El servicio que le dice a la aplicación de dónde bajar su actualización **construía la "
       "dirección a partir de la pregunta**, no de una configuración.",
       "Consultado desde el propio servidor respondía `http://localhost/…`, y desde internet "
       "devolvía `http://` en vez de `https://`: el cortafuegos institucional atiende el cifrado "
       "y entrega la petición en claro, así que el sistema creía que nadie usaba HTTPS."],
      tipo="problema")
panel(s, ML + W * 0.545, 1.92, W * 0.455, 3.15, "La solución",
      ["La dirección pública pasó a ser **una configuración explícita**, declarada donde se "
       "despliega, y ya no se deduce de nada.",
       "Si algún día faltara, el sistema la reconstruye respetando lo que anuncie el cortafuegos. "
       "Quedaron **cinco pruebas automáticas** que fijan el comportamiento en los tres escenarios."],
      tipo="solucion")
remate(s, "Es el mismo engaño del intermediario que ya nos había costado los registros de "
          "auditoría: lo que le llega al sistema no es lo que el usuario envió.", 5.35, 0.95)
pie(s, DS, 10)

# 11 · Dificultad 4 · la cadena del FTP
s = nueva(prs)
eyebrow(s, "Dificultad 4", "Caso 14512 · Medido de nuevo hoy, 4 de septiembre")
titulo(s, "Veinte noches seguidas sin cargar una sola encuesta", size=29)
panel(s, ML, 1.92, W * 0.455, 3.25, "El problema — sigue abierto",
      ["El proceso que sube al sistema las encuestas capturadas sin conexión corre todas las "
       "noches a las 18:20. **Su última carga exitosa fue el 14 de agosto.**",
       "Desde entonces falla **todas las noches, en un segundo**, siempre con el mismo error: "
       "perdió el acceso a su carpeta de trabajo tras un reinicio. Anoche volvió a fallar."],
      tipo="problema", chip=("20 noches", GRAVE, GRAVE_BG))
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.25, "Qué hace falta — y de quién depende",
      ["**Restablecer la carpeta** en el servidor donde vive el proceso. No es un cambio de "
       "programa: es un permiso de acceso. ~Administrador de ese servidor.~",
       "**Suspender la eliminación diaria** de los archivos del FTP. Hoy se borran a las 24 "
       "horas, sin comprobar que hayan entrado. ~Operación.~",
       "Ambas se escalaron por escrito el 1 de septiembre."],
      tipo="solucion")
nota(s, "**Medición de hoy:** 30 ejecuciones en los últimos 30 días — 10 correctas hasta el 14 de "
        "agosto y **20 fallidas seguidas** desde entonces. Cada noche que pasa se borra la "
        "captura de ese día antes de que alcance a cargarse.", 5.42, W * 0.94)
pie(s, DS, 11)

# 12 · Dificultad 5 · los correos del equipo
s = nueva(prs)
eyebrow(s, "Dificultad 5", "Capacitación · Bloquea la medición, no la jornada")
titulo(s, "El cuestionario identifica por correo, y seis del equipo no lo tienen", size=27)
panel(s, ML, 1.92, W * 0.48, 3.05, "El problema",
      ["Seis de los siete integrantes del equipo figuran en el sistema con direcciones del "
       "ambiente de desarrollo —`@srni.dev`, `@srni.local`—.",
       "El correo es lo que empareja el pre-test con el post-test. Sin el institucional, **la "
       "jornada se puede dictar pero no se puede medir la ganancia** de esas personas."],
      tipo="problema")
panel(s, ML + W * 0.52, 1.92, W * 0.48, 3.05, "La solución",
      ["Solicitar la creación o corrección de las seis cuentas con dominio "
       "`@unidadvictimas.gov.co` antes del **10 de septiembre**.",
       "Es la única condición de la Sesión 1 que no está en nuestras manos, y es de trámite: no "
       "bloquea nada más del proyecto."],
      tipo="solucion")
remate(s, "Se detectó al construir el cuestionario, no el día de la jornada. Corregirlo ahora "
          "cuesta un correo; el martes habría costado la medición completa.", 5.25, 0.95)
pie(s, DS, 12)

# 13 · Verificación
s = nueva(prs)
eyebrow(s, "Verificación", "Cómo se comprobó cada cosa")
titulo(s, "Nada de esto es una estimación")
tf = caja(s, ML, 1.78, W * 0.9, 0.6)
parrafo(tf, "Cada cifra de esta presentación tiene detrás una ejecución o una consulta, hecha en "
            "el periodo y repetible por un tercero.", 13.5, INK2, BODY, primero=True,
        space_after=0, line_spacing=1.3)
lista(s, ML, 2.55, W * 0.95,
      ["**1.185 pruebas automáticas** — batería completa del backend y de la aplicación móvil. El "
       "panel web no se pudo correr desde nuestro entorno: es una comprobación no realizada, no "
       "un fallo, y así se reporta.",
       "**Las cuatro direcciones publicadas** — verificadas contra el dominio institucional "
       "después del despliegue, no antes.",
       "**Las 10 preguntas** — comparadas una a una, en orden y opciones, entre el anexo impreso "
       "y lo que sirve el servidor.",
       "**Las veinte noches del proceso de carga** — leídas hoy del registro de ejecuciones del "
       "propio motor de base de datos, con la fecha y el error de cada corrida."],
      size=12.5, gap=0.12)
nota(s, "Las consultas son de solo lectura. No se modificó ningún dato, procedimiento ni trabajo "
        "programado del sistema legacy.", 5.72, W * 0.94)
pie(s, SEM, 13)

# 14 · Lo que sigue
s = nueva(prs)
eyebrow(s, "Lo que sigue", "Fechas y decisiones")
titulo(s, "La capacitación queda el 10, el 15 y el 18 de septiembre", size=28)
fila(s, 1.78, [("Acción", 0.55, PP_ALIGN.LEFT), ("De quién depende", 0.45, PP_ALIGN.LEFT)],
     alto=0.3, cabecera=True)
SIG = [("**Restablecer la carpeta de trabajo** del proceso de carga",
        "Administrador de ese servidor", "20 noches", CHIP_GRAVE),
       ("**Suspender la eliminación diaria** de archivos del FTP",
        "Operación", "Hoy", CHIP_GRAVE),
       ("**Sesión 1** · equipo de la Subdirección", "Equipo", "Jueves 10", CHIP_OK),
       ("**Sesión 2** · enlaces territoriales, Grupo A (16)",
        "Equipo · convocatoria", "Martes 15", CHIP_OK),
       ("**Sesión 3** · enlaces territoriales, Grupo B (14)",
        "Equipo · convocatoria", "Viernes 18", CHIP_OK),
       ("Correos institucionales de seis integrantes del equipo",
        "Administración", "Antes del 10", CHIP_PEND),
       ("Canal de soporte interno para los encuestadores",
        "Subdirección", "Bloquea el manual", CHIP_PEND)]
for i, (a, b, c, ch) in enumerate(SIG):
    fila(s, 2.22 + i * 0.55, [(a, 0.55, PP_ALIGN.LEFT), (b, 0.27, PP_ALIGN.LEFT),
                              ("", 0.18, PP_ALIGN.LEFT)], alto=0.47,
         chip=(c, ch[0], ch[1]))
nota(s, "Las dos primeras filas son las urgentes y ninguna es técnica nuestra: una es un permiso "
        "de carpeta y la otra es una instrucción a la operación. Llevan escaladas desde el 1 de "
        "septiembre.", 6.18, W * 0.94)
pie(s, SEM, 14)

# 15 · Cierre
s = nueva(prs)
eyebrow(s, "Cierre", "28 ago → 4 sep 2026")
titulo(s, "La capacitación quedó lista, y con forma de medirla")
tiles(s, [("1.185", "pruebas en verde, 0 fallos"),
          ("4", "recursos publicados: app, manual, pre-test y post-test"),
          ("5", "dificultades resueltas o escaladas, cuatro halladas por nosotros"),
          ("20", "noches del proceso de carga sin operar, escaladas y sin respuesta")],
      2.05, h=1.55)
remate(s, "Cuatro de las cinco dificultades aparecieron porque revisamos nuestro propio trabajo "
          "antes de entregarlo. La quinta lleva veinte noches esperando una decisión que no es "
          "nuestra.", 4.05, 1.15)
nota(s, "Alexandra López · Dirección técnica RNI   ·   Javier Aguilar · Arquitectura y desarrollo "
        "  ·   Brandon Niño · Panel de Control   ·   Jorge Cardona · Calidad   ·   Karen Serna · "
        "Documentación   ·   Nixon · Seguridad   ·   Oscar Manosalva · Supervisión   ·   "
        "PRY-0662064", 5.55, W * 0.94)
pie(s, SEM, 15)

# ─── Guardar ────────────────────────────────────────────────────────────────
import os
DEST = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'pptx',
                    'presentacion_avance_28ago-4sep.pptx')
prs.save(os.path.normpath(DEST))
print("PPTX generado:", os.path.normpath(DEST))
print("diapositivas:", len(prs.slides.__iter__.__self__._sldIdLst))
