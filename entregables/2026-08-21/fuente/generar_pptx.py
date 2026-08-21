# -*- coding: utf-8 -*-
"""
Genera la presentacion de avance semanal (14-ago -> 21-ago-2026) en .pptx.

Replica el diseno del HTML de esta misma carpeta: marca institucional #ffcc03,
tinta calida, geometria recta. La estructura es avance / problema / solucion,
con cada problema y su solucion en la MISMA diapositiva.

    python generar_pptx.py

Salida: ../pptx/presentacion_avance_14-21-ago.pptx
"""
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ─── Paleta (identica al HTML, tema claro) ──────────────────────────────────
INK        = RGBColor(0x17, 0x16, 0x0F)
INK2       = RGBColor(0x4A, 0x47, 0x3B)
INK3       = RGBColor(0x7C, 0x78, 0x69)
SURFACE    = RGBColor(0xFB, 0xFA, 0xF6)
PANEL      = RGBColor(0xF2, 0xEF, 0xE4)
LINE       = RGBColor(0xDC, 0xD8, 0xC7)
LINE_SOFT  = RGBColor(0xEB, 0xE7, 0xD9)
MARCA      = RGBColor(0xFF, 0xCC, 0x03)
MARCA_INK  = RGBColor(0x7A, 0x5C, 0x00)
OK         = RGBColor(0x1F, 0x6B, 0x3A)
OK_BG      = RGBColor(0xE6, 0xF0, 0xE6)
ALERTA     = RGBColor(0x9C, 0x5A, 0x00)
ALERTA_BG  = RGBColor(0xF7, 0xEC, 0xD7)
GRAVE      = RGBColor(0xA8, 0x1F, 0x1A)
GRAVE_BG   = RGBColor(0xF8, 0xE5, 0xE2)
MARK_ANTES = RGBColor(0xC9, 0xC4, 0xB0)

DISPLAY = "Segoe UI"
BODY    = "Segoe UI"
MONO    = "Consolas"

# ─── Geometria (pulgadas) ───────────────────────────────────────────────────
SW, SH   = 13.333, 7.5
ML       = 0.78
W        = SW - ML * 2
FILETE_H = 0.06
EYEBROW_Y = 0.46
TITULO_Y  = 0.82
CUERPO_Y  = 1.82
PIE_Y     = 6.86
PIE_LINEA = 6.80

TOTAL = 15


# ─── Utilidades ─────────────────────────────────────────────────────────────

def spc(run, centesimas):
    """Espaciado entre caracteres (no expuesto por python-pptx)."""
    run.font._rPr.set('spc', str(int(centesimas)))


def rect(slide, l, t, w, h, color, linea=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if linea is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = linea
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def caja(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


TROZO = re.compile(r'(\*\*.+?\*\*|~.+?~)')


def escribir(p, texto, size, color, font=BODY, bold=False, fuerte=None):
    """Escribe texto con **negrita** y ~acento~ como runs."""
    fuerte = fuerte or INK
    for trozo in TROZO.split(texto):
        if not trozo:
            continue
        r = p.add_run()
        if trozo.startswith('**') and trozo.endswith('**'):
            r.text = trozo[2:-2]
            r.font.bold = True
            r.font.color.rgb = fuerte
        elif trozo.startswith('~') and trozo.endswith('~'):
            r.text = trozo[1:-1]
            r.font.bold = True
            r.font.color.rgb = MARCA_INK
        else:
            r.text = trozo
            r.font.bold = bold
            r.font.color.rgb = color
        r.font.size = Pt(size)
        r.font.name = font
    return p


def parrafo(tf, texto, size, color, font=BODY, bold=False, primero=False,
            space_after=6, line_spacing=1.28, align=None, fuerte=None):
    p = tf.paragraphs[0] if primero else tf.add_paragraph()
    escribir(p, texto, size, color, font, bold, fuerte)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    return p


# ─── Componentes ────────────────────────────────────────────────────────────

def nueva(prs, fondo=SURFACE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = fondo
    rect(s, 0, 0, SW, FILETE_H, MARCA)
    return s


def eyebrow(s, etiqueta, resto):
    tf = caja(s, ML, EYEBROW_Y, W, 0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = etiqueta.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.bold = True
    r.font.color.rgb = MARCA_INK; spc(r, 140)
    r = p.add_run(); r.text = "  ·  " + resto.upper()
    r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 140)


def titulo(s, texto, size=30, y=TITULO_Y, h=0.9):
    tf = caja(s, ML, y, W, h)
    parrafo(tf, texto, size, INK, DISPLAY, bold=True, primero=True,
            space_after=0, line_spacing=1.06)


def pie(s, izq, n):
    rect(s, ML, PIE_LINEA, W, 0.008, LINE)
    tf = caja(s, ML, PIE_Y, W * 0.7, 0.25)
    p = parrafo(tf, izq.upper(), 9, INK3, MONO, primero=True, space_after=0, line_spacing=1)
    for r in p.runs:
        spc(r, 120)
    tf = caja(s, ML + W * 0.7, PIE_Y, W * 0.3, 0.25)
    p = parrafo(tf, "%02d / %d" % (n, TOTAL), 9, INK3, MONO, primero=True,
                space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)
    for r in p.runs:
        spc(r, 120)


def tiles(s, items, y, h=1.42):
    """items: [(cifra, etiqueta), ...]"""
    n = len(items)
    gap = 0.22
    w = (W - gap * (n - 1)) / n
    for i, (cifra, etiqueta) in enumerate(items):
        l = ML + i * (w + gap)
        rect(s, l, y, w, h, PANEL)
        rect(s, l, y, w, 0.055, MARCA)
        tf = caja(s, l + 0.26, y + 0.26, w - 0.5, h - 0.4)
        parrafo(tf, cifra, 28, INK, MONO, bold=True, primero=True,
                space_after=5, line_spacing=1)
        parrafo(tf, etiqueta, 10.5, INK2, BODY, space_after=0, line_spacing=1.22)


def deltas(s, items, y, alto=0.62, gap=0.13):
    """items: [(que, antes, ahora), ...]"""
    for i, (que, antes, ahora) in enumerate(items):
        t = y + i * (alto + gap)
        rect(s, ML, t, W, alto, PANEL)
        rect(s, ML, t, 0.07, alto, MARCA)
        tf = caja(s, ML + 0.32, t + 0.15, W * 0.52, alto - 0.2)
        parrafo(tf, que, 13, INK2, BODY, primero=True, space_after=0, line_spacing=1.1)
        tf = caja(s, ML + W * 0.55, t + 0.16, W * 0.2, alto - 0.2)
        parrafo(tf, antes, 13, INK3, MONO, primero=True, space_after=0,
                line_spacing=1.1, align=PP_ALIGN.RIGHT)
        tf = caja(s, ML + W * 0.755, t + 0.15, 0.3, alto - 0.2)
        parrafo(tf, "→", 13, MARCA_INK, BODY, bold=True, primero=True,
                space_after=0, line_spacing=1.1, align=PP_ALIGN.CENTER)
        tf = caja(s, ML + W * 0.79, t + 0.12, W * 0.21 - 0.1, alto - 0.2)
        parrafo(tf, ahora, 16, INK, MONO, bold=True, primero=True, space_after=0,
                line_spacing=1.05, align=PP_ALIGN.RIGHT)


def panel(s, l, t, w, h, rotulo, parrafos, tipo="neutro", chip=None):
    fondo, barra = PANEL, LINE
    if tipo == "problema":
        fondo, barra = GRAVE_BG, GRAVE
    elif tipo == "solucion":
        fondo, barra = PANEL, MARCA
    rect(s, l, t, w, h, fondo)
    rect(s, l, t, 0.07, h, barra)

    tf = caja(s, l + 0.32, t + 0.24, w - 0.6, 0.26)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = rotulo.upper()
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 130)
    p.space_after = Pt(0); p.line_spacing = 1
    if chip:
        texto, color, bg = chip
        cw = 0.15 + len(texto) * 0.072
        cl = l + 0.32 + 1.45
        rect(s, cl, t + 0.2, cw, 0.26, bg, linea=color)
        ctf = caja(s, cl, t + 0.245, cw, 0.22)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 80)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1

    tf = caja(s, l + 0.32, t + 0.62, w - 0.62, h - 0.85)
    for i, txt in enumerate(parrafos):
        parrafo(tf, txt, 12.5, INK2, BODY, primero=(i == 0), space_after=9, line_spacing=1.26)


def vineta(p):
    """Vineta cuadrada de marca, alineada por PowerPoint a la primera linea."""
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '182880')
    pPr.set('indent', '-182880')
    sucesores = ('a:buSzTx', 'a:buSzPct', 'a:buSzPts', 'a:buFontTx', 'a:buFont',
                 'a:buNone', 'a:buAutoNum', 'a:buChar', 'a:tabLst', 'a:defRPr', 'a:extLst')
    buClr = pPr.makeelement(qn('a:buClr'), {})
    buClr.append(pPr.makeelement(qn('a:srgbClr'), {'val': 'FFCC03'}))
    pPr.insert_element_before(buClr, *sucesores)
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    pPr.insert_element_before(buFont, 'a:buNone', 'a:buAutoNum', 'a:buChar',
                              'a:tabLst', 'a:defRPr', 'a:extLst')
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    pPr.insert_element_before(buChar, 'a:tabLst', 'a:defRPr', 'a:extLst')


def lista(s, l, t, w, items, size=12.5, gap=0.1):
    tf = caja(s, l, t, w, 3.4)
    for i, txt in enumerate(items):
        p = parrafo(tf, txt, size, INK2, BODY, primero=(i == 0),
                    space_after=13, line_spacing=1.26)
        vineta(p)


def nota(s, texto, y, w=None):
    tf = caja(s, ML, y, w or W, 0.6)
    parrafo(tf, texto, 10.5, INK3, BODY, primero=True, space_after=0, line_spacing=1.3)


def remate(s, texto, y, h=0.9):
    rect(s, ML, y, 0.07, h, MARCA)
    tf = caja(s, ML + 0.32, y + 0.06, W - 0.4, h)
    parrafo(tf, texto, 17, INK, DISPLAY, primero=True, space_after=0, line_spacing=1.22)


def fila(s, y, celdas, alto=0.62, cabecera=False, sub=None, chip=None):
    """celdas: [(texto, ancho_relativo, align), ...] sobre el ancho util."""
    x = ML
    for texto, wr, align in celdas:
        cw = W * wr
        tf = caja(s, x, y + (0.02 if cabecera else 0.06), cw - 0.2, alto)
        if cabecera:
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = texto.upper()
            r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 110)
            p.alignment = align; p.space_after = Pt(0); p.line_spacing = 1
        else:
            parrafo(tf, texto, 12.5, INK2, BODY, primero=True, space_after=0,
                    line_spacing=1.2, align=align)
        x += cw
    if sub:
        tf = caja(s, ML, y + 0.32, W * 0.5, 0.4)
        parrafo(tf, sub, 10, INK3, BODY, primero=True, space_after=0, line_spacing=1.2)
    if chip:
        texto, color, bg = chip
        cw = 0.18 + len(texto) * 0.075
        rect(s, ML + W - cw, y + 0.05, cw, 0.28, bg, linea=color)
        ctf = caja(s, ML + W - cw, y + 0.1, cw, 0.24)
        p = ctf.paragraphs[0]
        r = p.add_run(); r.text = texto.upper()
        r.font.size = Pt(8.5); r.font.name = MONO; r.font.color.rgb = color; r.font.bold = True
        spc(r, 70)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.line_spacing = 1
    rect(s, ML, y + alto, W, 0.008, INK if cabecera else LINE_SOFT)


# ─── Las 15 diapositivas ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)

CHIP_OK   = (OK, OK_BG)
CHIP_PEND = (ALERTA, ALERTA_BG)

# 01 · Portada
s = nueva(prs)
tf = caja(s, ML, 0.85, W, 0.3)
p = parrafo(tf, "Unidad para las Víctimas  ·  PRY-0662064".upper(), 10.5, INK3, MONO,
            primero=True, space_after=0, line_spacing=1)
for r in p.runs:
    spc(r, 160)
tf = caja(s, ML, 1.95, W * 0.94, 2.6)
parrafo(tf, "Calidad revisó la aplicación, y el equipo cerró los siete hallazgos "
            "y otros doce que nadie había visto.", 38, INK, DISPLAY, bold=True,
        primero=True, space_after=0, line_spacing=1.03)
rect(s, ML, 4.86, 4.6, 0.035, MARCA)
tf = caja(s, ML, 5.06, W, 0.35)
parrafo(tf, "Avance semanal  ·  14 ago → 21 ago 2026", 14, INK, MONO, primero=True,
        space_after=0, line_spacing=1)
tf = caja(s, ML, 6.55, W * 0.6, 0.3)
parrafo(tf, "SICAV Móvil  ·  Sistema de Caracterización de Víctimas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML + W * 0.5, 6.55, W * 0.5, 0.3)
parrafo(tf, "26 cambios versionados en 5 días de trabajo", 10.5, INK3, MONO, primero=True,
        space_after=0, line_spacing=1, align=PP_ALIGN.RIGHT)

# 02 · El periodo en cifras
s = nueva(prs)
eyebrow(s, "Resumen", "El periodo en cifras")
titulo(s, "La semana del informe de calidad")
tf = caja(s, ML, 1.75, W * 0.86, 0.7)
parrafo(tf, "El informe **IGED-QA-C003** llegó con trece hallazgos. Siete eran defectos por "
            "corregir; los otros seis el propio informe los daba por cumplidos.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.32)
tiles(s, [("13 / 13", "hallazgos de calidad atendidos"),
          ("19", "defectos corregidos: 7 del informe y 12 encontrados por el equipo"),
          ("1.098", "pruebas automáticas en verde ~(+91)~"),
          ("3", "versiones de la aplicación publicadas")], 2.72, h=1.62)
remate(s, "Doce de los diecinueve defectos no los reportó nadie: aparecieron al abrir el "
          "trabajo sin conexión y mirar qué había detrás.", 4.72, 1.1)
pie(s, "SICAV Móvil · Avance semanal", 2)

# 03 · Divisor avances
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.55, W, 1.5)
parrafo(tf, "01", 76, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.15, W, 1.0)
parrafo(tf, "Avances", 44, INK, DISPLAY, bold=True, primero=True, space_after=0,
        line_spacing=1.02)
tf = caja(s, ML, 4.35, W * 0.72, 1.1)
parrafo(tf, "Cuatro frentes cerrados: quién autoriza una excepción, el trabajo sin conexión, "
            "el informe de calidad completo y lo que ya quedó corriendo en producción.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.35)
pie(s, "SICAV Móvil · Avance semanal", 3)

# 04 · Avance 1 · la excepción cambió de manos
s = nueva(prs)
eyebrow(s, "Avance 1", "Quién autoriza una excepción")
titulo(s, "Autorizar dejó de ser tarea de quien está en el barrio")
deltas(s, [("Quién decide saltarse el control de vigencia", "la encuestadora", "coordinación"),
           ("Qué se le pedía en campo", "foto del fallo o la tutela", "nada"),
           ("Documentos por autorización", "uno a uno", "hasta 200")], 2.15, alto=0.72, gap=0.16)
nota(s, "La habilitación viaja al teléfono en la precarga de la jornada y queda guardada, no en "
        "memoria: **funciona sin señal**. El celular ya no decide, solo consume lo que "
        "coordinación autorizó.", 5.0, W * 0.9)
pie(s, "Avances", 4)

# 05 · Avance 2 · sin señal
s = nueva(prs)
eyebrow(s, "Avance 2", "El trabajo sin conexión")
titulo(s, "En modo avión la aplicación había quedado muda")
fila(s, 1.95, [("Pantalla", 0.18, PP_ALIGN.LEFT), ("Antes, sin señal", 0.41, PP_ALIGN.LEFT),
               ("Ahora", 0.41, PP_ALIGN.LEFT)], alto=0.34, cabecera=True)
for i, (a, b, c) in enumerate([
        ("**Hogares**", "«Sin hogares», aunque hubiera hogares conformados en el teléfono",
         "Los muestra desde la base local, con su estado real"),
        ("**Encuestas**", "«No se pudo cargar las sesiones»",
         "Muestra las entrevistas guardadas, con aviso de que falta enviarlas"),
        ("**Detalle de sesión**", "Una entrevista en curso no se podía abrir",
         "Abre sus respuestas y permite continuarla")]):
    fila(s, 2.45 + i * 0.86, [(a, 0.18, PP_ALIGN.LEFT), (b, 0.41, PP_ALIGN.LEFT),
                              (c, 0.41, PP_ALIGN.LEFT)], alto=0.78)
nota(s, "El equipo también construyó una barra de progreso propia: la anterior se salía de su "
        "tarjeta cuando el porcentaje pasaba de 100. Ese era el hallazgo APK-006.", 5.4, W * 0.9)
pie(s, "Avances", 5)

# 06 · Avance 3 · el informe de calidad
s = nueva(prs)
eyebrow(s, "Avance 3", "El informe de calidad")
titulo(s, "Los siete defectos, atendidos de punta a punta", size=28)
fila(s, 1.85, [("Hallazgo", 0.13, PP_ALIGN.LEFT), ("Qué reportó calidad", 0.62, PP_ALIGN.LEFT),
               ("Estado", 0.25, PP_ALIGN.LEFT)], alto=0.3, cabecera=True)
FILAS = [("**APK-001**", "La ruta de excepción no avanzaba", "Resuelto", CHIP_OK),
         ("**APK-002**", "«No se pudo registrar» al conformar hogar", "Resuelto", CHIP_OK),
         ("**APK-003**", "El modo sin conexión no funcionaba", "Resuelto", CHIP_OK),
         ("**APK-004**", "No se podía quitar ni corregir un integrante", "Quitar sí · falta corregir", CHIP_PEND),
         ("**APK-005**", "Sesión «Completada» con la barra en 0 %", "Resuelto", CHIP_OK),
         ("**APK-006**", "Las barras se salían de la tarjeta", "Resuelto", CHIP_OK),
         ("**APK-007**", "«No habilitado» sin el nombre de la persona", "Resuelto", CHIP_OK)]
for i, (cod, txt, est, (col, bg)) in enumerate(FILAS):
    fila(s, 2.3 + i * 0.46, [(cod, 0.13, PP_ALIGN.LEFT), (txt, 0.62, PP_ALIGN.LEFT),
                             ("", 0.25, PP_ALIGN.LEFT)], alto=0.4,
         chip=(est, col, bg))
nota(s, "Del **APK-008 al APK-013** —autenticación, alerta de vigencia, exactitud de los datos, "
        "captura por Ruta General, validación de campos y diseño del mecanismo de excepción— el "
        "propio informe los da por cumplidos. **Dos de los siete no eran lo que parecían.**",
     5.62, W * 0.94)
pie(s, "Avances", 6)

# 07 · Avance 4 · en producción
s = nueva(prs)
eyebrow(s, "Avance 4", "Lo que quedó corriendo")
titulo(s, "Los tres componentes, desplegados y verificados")
tiles(s, [("Servidor", "imagen nueva, cuatro servicios recreados, sin migraciones pendientes"),
          ("Panel web", "construido y sirviendo la versión nueva"),
          ("1.2.2", "aplicación publicada y descargable por el dominio institucional")],
      1.9, h=1.5)
lista(s, ML, 3.62, W * 0.94, [
    "Las tres rutas responden por **caracterizacion.unidadvictimas.gov.co**: panel, servicios y descarga de la aplicación.",
    "El código nuevo se comprobó **corriendo en el servidor**, no solo desplegado: se evaluaron reglas reales del instrumento contra el motor en producción.",
    "El código QR de descarga **no cambia entre versiones**: quien ya lo tenga impreso baja la versión nueva sin hacer nada.",
    "La versión anterior queda respaldada en el servidor: volver atrás es copiar un archivo."], size=12.5)
pie(s, "Avances", 7)

# 08 · Divisor dificultades
s = nueva(prs, PANEL)
tf = caja(s, ML, 1.55, W, 1.5)
parrafo(tf, "02", 76, MARCA, MONO, bold=True, primero=True, space_after=0, line_spacing=0.9)
tf = caja(s, ML, 3.15, W, 1.0)
parrafo(tf, "Dificultades y soluciones", 40, INK, DISPLAY, bold=True, primero=True,
        space_after=0, line_spacing=1.02)
tf = caja(s, ML, 4.35, W * 0.76, 1.2)
parrafo(tf, "Cuatro obstáculos del periodo. Dos son hallazgos de calidad cuya causa resultó ser "
            "otra que la reportada, y dos son de fondo: el trabajo que desaparecía sin señal y "
            "un número que llevaba tiempo mintiendo.",
        14, INK2, BODY, primero=True, space_after=0, line_spacing=1.35)
pie(s, "SICAV Móvil · Avance semanal", 8)

# 09 · Dificultad 1 · la foto
s = nueva(prs)
eyebrow(s, "Dificultad 1", "APK-001 · La foto que nadie en campo podía tener")
titulo(s, "El botón no era el problema")
panel(s, ML, 1.92, W * 0.455, 3.42, "El problema",
      ["Calidad reportó que la ruta de excepción **no dejaba avanzar**. Al desarmarlo aparecieron "
       "dos causas, y ninguna era la reportada.",
       "La de negocio: el botón le pedía a la encuestadora una **foto del fallo o la tutela**, un "
       "documento que ella no tiene — llega por canal institucional al nivel central.",
       "La técnica: el registro de auditoría se escribía con un nombre de campo equivocado, así "
       "que **toda llamada moría en error 500** antes de responder. Ese camino nunca funcionó, ni "
       "antes del informe."],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.42, "La solución",
      ["No se parchó el botón: **se cambió quién autoriza**. La excepción la otorga coordinación "
       "desde el panel, con radicado y motivo, y en lotes de hasta 200 documentos.",
       "El celular solo la consume. La ve al buscar a la persona o en la precarga de la jornada, y "
       "queda guardada en el teléfono: ~funciona sin señal~.",
       "El camino viejo responde ahora **«ya no existe»** y no «no encontrado»: un «no encontrado» "
       "lo lee la versión anterior de la aplicación como falta de red, y mandaría a la encuestadora "
       "a buscar señal por un camino que no va a volver."],
      tipo="solucion")
nota(s, "Que nadie lo hubiera notado tiene explicación: hasta hoy ninguna encuestadora ha entrado "
        "al sistema. El informe de calidad fue la primera vez que alguien recorrió ese camino.",
     5.55, W * 0.94)
pie(s, "Dificultades y soluciones", 9)

# 10 · Dificultad 2 · el error intermitente
s = nueva(prs)
eyebrow(s, "Dificultad 2", "APK-002 · El error «intermitente»")
titulo(s, "Lo intermitente era el mensaje, no el error")
panel(s, ML, 1.92, W * 0.455, 3.42, "El problema",
      ["Con algunos documentos, conformar el hogar mostraba «No se pudo registrar. Revisa la "
       "conexión». Con otros funcionaba. **Parecía un defecto intermitente.**",
       "Tres situaciones distintas —sin red, la persona ya pertenece al hogar de otro encuestador, "
       "y respuesta vacía del servidor— se veían **exactamente iguales**.",
       "En uno de esos casos, la pantalla llegaba a mostrarle a la encuestadora literalmente la "
       "palabra «null»."],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.42, "La solución",
      ["Un intérprete de errores que separa los tres casos, porque cada uno pide una acción "
       "distinta de la encuestadora.",
       "**Sin red** → «Su trabajo no se pierde: vuelva a intentarlo cuando tenga señal».",
       "**El servidor respondió** → su texto, que ya viene redactado para campo.",
       "**Cualquier otra cosa** → mensaje claro más el código, y el detalle técnico va al reporte "
       "de errores, no a la pantalla."],
      tipo="solucion")
nota(s, "Los dos documentos que cita el informe caían en «hogar de otro encuestador»: eso ahora se "
        "lee. No había un defecto intermitente que perseguir — faltaba que soporte pudiera "
        "diagnosticar sin pedir capturas de pantalla.", 5.55, W * 0.94)
pie(s, "Dificultades y soluciones", 10)

# 11 · Dificultad 3 · el trabajo que desaparecía
s = nueva(prs)
eyebrow(s, "Dificultad 3", "APK-003 · El trabajo que desaparecía")
titulo(s, "Se abrió la puerta del trabajo sin conexión y detrás había doce cosas rotas",
       size=26)
panel(s, ML, 1.92, W * 0.455, 3.42, "El problema",
      ["Mientras las pantallas no mostraban nada sin señal, **nadie podía ver lo que estaba mal "
       "detrás**. Al hacerlas funcionar aparecieron doce defectos, todos de la misma familia.",
       "Los tres peores: una entrevista **se partía en dos** y al volver a entrar se abría un "
       "formulario en blanco; la lista de hogares **borraba copias locales** porque trataba una "
       "página de 20 resultados como si fuera el listado completo; y cerrar dos veces la misma "
       "encuesta dejaba el teléfono en un estado del que no salía — y con eso, **sin borrar los "
       "datos personales al cerrar sesión**."],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.42, "La solución",
      ["Cuatro revisiones sucesivas: **3, 5, 2 y 2 defectos**. Las dos últimas buscaban "
       "específicamente si los arreglos anteriores habían roto algo, y las dos encontraron algo.",
       "Se separó «cerrada en el teléfono» de «cerrada en el servidor», que era una sola palabra "
       "para dos hechos distintos. Va con una ~migración que repara los teléfonos que ya vienen de "
       "campo~ con filas mal clasificadas.",
       "Y se escribieron **pruebas que fallan si el defecto vuelve**: se comprobó revirtiendo cada "
       "arreglo a mano."],
      tipo="solucion")
nota(s, "Todos comparten una raíz: se mostraba lo que hay en la base del teléfono, pero las "
        "condiciones que deciden qué mostrar y qué borrar estaban escritas pensando en que hay red.",
     5.55, W * 0.94)
pie(s, "Dificultades y soluciones", 11)

# 12 · Dificultad 4 · la barra en 0 %
s = nueva(prs)
eyebrow(s, "Dificultad 4", "APK-005 · El número que mentía")
titulo(s, "Al porcentaje le faltaba una palabra: «visibles»")
panel(s, ML, 1.92, W * 0.455, 3.42, "El problema",
      ["Entrevistas marcadas «Completada» con la barra en **0 %**. El primer arreglo fue mostrar "
       "100 % cuando el estado fuera «Completada», y **el equipo lo retiró**: tapaba el síntoma y "
       "además mentía, porque una entrevista interrumpida a mitad —la víctima se retiró— también "
       "se veía al 100 %.",
       "El panel web nunca aplicó ese ajuste, así que la misma entrevista se veía distinta según "
       "quién la mirara.",
       "La causa estaba en el servidor: el cálculo dividía por **todas** las preguntas "
       "obligatorias, incluidas las que las reglas del formulario mantienen ocultas y que **nadie "
       "puede responder nunca**."],
      tipo="problema")
panel(s, ML + W * 0.495, 1.92, W * 0.505, 3.42, "La solución",
      ["El denominador pasó a ser ~las obligatorias que la persona realmente ve~, evaluando las "
       "reglas con sus datos: edad, sexo, pertenencia étnica e inclusión en el registro.",
       "Una pregunta puede aplicarle a un integrante del hogar y no a otro, y ahora se cuenta así.",
       "El motor de reglas quedó en **un solo lugar**. Antes cada parte del sistema tenía su copia; "
       "hoy la captura en el celular, el tablero del celular y el servidor deciden con el mismo "
       "criterio."],
      tipo="solucion")
nota(s, "Se reabrió el hallazgo a propósito para atacar la causa. El porcentaje es lo que la "
        "supervisión mira para saber si una entrevista terminó: mientras el móvil dijera una cosa y "
        "el panel otra, ese número no servía para decidir.", 5.55, W * 0.94)
pie(s, "Dificultades y soluciones", 12)

# 13 · Verificación
s = nueva(prs)
eyebrow(s, "Verificación", "Cómo sabemos que está bien")
titulo(s, "Noventa y una pruebas nuevas, y cada una comprobada")
for i, (rot, ancho, val, color) in enumerate(
        [("Servidor", 0.935, "883 → 944", MARCA),
         ("Aplicación", 0.821, "115 → 140", MARCA),
         ("Panel web", 0.643, "9 → 14", MARCA)]):
    yb = 1.95 + i * 0.5
    tf = caja(s, ML, yb + 0.06, 1.1, 0.3)
    parrafo(tf, rot, 11, INK2, MONO, primero=True, space_after=0, line_spacing=1)
    pista_l, pista_w = ML + 1.2, W - 2.6
    rect(s, pista_l, yb, pista_w, 0.34, LINE_SOFT)
    rect(s, pista_l, yb, pista_w * ancho, 0.34, color)
    tf = caja(s, ML + W - 1.35, yb + 0.05, 1.35, 0.3)
    parrafo(tf, val, 12, INK, MONO, bold=True, primero=True, space_after=0,
            line_spacing=1, align=PP_ALIGN.RIGHT)
nota(s, "Total de las tres: **1.007 → 1.098 pruebas** en verde. Medido corriendo las suites en el "
        "punto de partida y hoy, no estimado.", 3.55, W * 0.94)
lista(s, ML, 4.15, W * 0.94, [
    "Las pruebas se comprobaron **por mutación**: se revirtió cada arreglo a mano para confirmar que la prueba falla cuando el defecto vuelve. Una prueba que no falla no es una prueba.",
    "Una de ellas ya se pagó sola: atajó un problema de rendimiento **antes de que saliera del equipo de desarrollo**, en el cálculo que corre cada vez que se guarda una respuesta.",
    "**Lo que estas pruebas no cubren:** son pruebas de la lógica, no de las pantallas. Quedan tres verificaciones en dispositivo real."], size=12)
pie(s, "SICAV Móvil · Avance semanal", 13)

# 14 · Lo que sigue
s = nueva(prs)
eyebrow(s, "Lo que sigue", "Y una decisión que necesitamos")
titulo(s, "Una sola cosa puede frenar el arranque, y no es técnica")
panel(s, ML, 1.92, W * 0.475, 3.5, "La decisión que llevamos a la supervisión",
      ["El rediseño movió la autorización de excepciones al nivel central. Hoy en producción hay "
       "**un coordinador, un supervisor y un administrador** con ese permiso.",
       "Para **1.158 encuestadoras**.",
       "Si no se define cuántas cuentas se habilitan y quiénes son, el flujo nuevo se traba el "
       "primer día de operación. **No es una decisión del equipo.**"],
      tipo="problema")
rect(s, ML + W * 0.515, 1.92, W * 0.485, 3.5, PANEL)
rect(s, ML + W * 0.515, 1.92, 0.07, 3.5, MARCA)
tf = caja(s, ML + W * 0.515 + 0.32, 2.16, W * 0.485 - 0.6, 0.26)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "LO QUE SIGUE DEL LADO DEL EQUIPO"
r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = INK3; spc(r, 130)
p.space_after = Pt(0); p.line_spacing = 1
lista(s, ML + W * 0.515 + 0.32, 2.52, W * 0.485 - 0.66, [
    "Verificar en dispositivo real los tres hallazgos que quedaron esperando el build nuevo",
    "La pantalla para corregir los datos de un integrante — cierra el APK-004",
    "Documentar hasta dónde llega el trabajo sin conexión y qué no cubre",
    "Llevar la regla de recaracterización al manual de usuario",
    "Un hilo suelto ya identificado: un envío fallido puede impedir la limpieza de la base local del teléfono con el uso prolongado"], size=11)
nota(s, "Estado del proyecto: estabilización previa al arranque. Esta semana se cerró el último "
        "bloqueo crítico reportado por calidad.", 5.62, W * 0.94)
pie(s, "SICAV Móvil · Avance semanal", 14)

# 15 · Cierre
s = nueva(prs)
tf = caja(s, ML, 0.85, W, 0.3)
p = parrafo(tf, "Equipo SICAV Móvil".upper(), 10.5, INK3, MONO, primero=True,
            space_after=0, line_spacing=1)
for r in p.runs:
    spc(r, 160)
tf = caja(s, ML, 2.0, W * 0.92, 1.8)
parrafo(tf, "Los siete hallazgos de calidad, cerrados. Y doce que nadie había reportado.",
        38, INK, DISPLAY, bold=True, primero=True, space_after=0, line_spacing=1.03)
rect(s, ML, 3.86, 4.2, 0.035, MARCA)
tf = caja(s, ML, 4.06, W, 0.35)
parrafo(tf, "Corte: 21 de agosto de 2026", 14, INK, MONO, primero=True, space_after=0,
        line_spacing=1)
lista(s, ML, 4.8, W * 0.47, [
    "**Javier Aguilar** — arquitectura, servidor, base de datos, aplicación móvil e infraestructura",
    "**Brando** — panel web y trabajo sin conexión de la aplicación"], size=12)
lista(s, ML + W * 0.53, 4.8, W * 0.47, [
    "**Jorge** — aseguramiento de calidad · informe IGED-QA-C003",
    "**Oscar** — supervisión del proyecto por parte de la Unidad"], size=12)
tf = caja(s, ML, 6.55, W * 0.62, 0.3)
parrafo(tf, "Cifras medidas contra el repositorio y el servidor, no estimadas", 10.5, INK3, MONO,
        primero=True, space_after=0, line_spacing=1)
tf = caja(s, ML + W * 0.5, 6.55, W * 0.5, 0.3)
parrafo(tf, "15 / 15", 10.5, INK3, MONO, primero=True, space_after=0,
        line_spacing=1, align=PP_ALIGN.RIGHT)

destino = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "pptx")
os.makedirs(destino, exist_ok=True)
ruta = os.path.normpath(os.path.join(destino, "presentacion_avance_14-21-ago.pptx"))
prs.save(ruta)
print("OK:", ruta, "| diapositivas:", len(prs.slides.__iter__.__self__._sldIdLst))
